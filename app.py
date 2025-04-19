from flask import Flask, request, send_file, render_template, jsonify, url_for, make_response
from flask_talisman import Talisman # Security Headers
from flask_wtf.csrf import CSRFProtect, CSRFError # CSRF Protection
from flask_limiter import Limiter # Rate Limiting
from flask_limiter.util import get_remote_address # Rate Limiting Helper
import os
import sys
import time
import subprocess
import logging
import glob # <--- THÊM IMPORT NÀY
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge # Better handling for large files
from pdf2docx import Converter # Cần cho bước chuyển đổi cuối
import tempfile
import PyPDF2
import shutil # Cần cho safe_remove, which
from pdf2image import convert_from_path, pdfinfo_from_path
from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image, UnidentifiedImageError
import zipfile
try:
    import magic # MIME Type Detection - Thư viện này có thể khó cài trên Windows
except ImportError:
    magic = None
    logging.warning("python-magic library not found. MIME type detection might be less reliable.")
from werkzeug.middleware.proxy_fix import ProxyFix

# === Basic Flask App Setup ===
app = Flask(__name__, template_folder='templates', static_folder='static')

# === Configuration ===
app.config['MAX_CONTENT_LENGTH'] = 101 * 1024 * 1024  # ~101MB limit server-side
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-secret-key-change-in-prod')
app.config['WTF_CSRF_SECRET_KEY'] = app.config['SECRET_KEY']
app.wsgi_app = ProxyFix(
    app.wsgi_app, x_for=1, x_proto=1, x_host=1
)

# === Logging ===
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - [%(funcName)s] - %(message)s',
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)

# === Security ===
csrf = CSRFProtect(app)
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["419 per day", "210 per hour", "30 per minute"],
    storage_uri="memory://",
    strategy="fixed-window"
)
csp = {
    'default-src': ['\'self\'', 'https://cdn.tailwindcss.com', 'https://fonts.googleapis.com', 'https://fonts.gstatic.com'],
    'style-src': ['\'self\'', '\'unsafe-inline\'', 'https://cdn.tailwindcss.com', 'https://fonts.googleapis.com'], # Allow inline styles for now
    'script-src': ['\'self\'', '\'unsafe-inline\'', 'https://cdn.tailwindcss.com'], # Allow inline scripts for now
    'font-src': ['\'self\'', 'https://fonts.gstatic.com'],
    'img-src': ['\'self\'', 'data:'],
    'form-action': '\'self\''
}
# Tạm thời tắt CSP nếu bạn đang gặp vấn đề với nó khi debug local
# talisman = Talisman(app, content_security_policy=None)
talisman = Talisman(
    app,
    content_security_policy=csp,
    force_https=False, # Set True in production behind a proxy handling HTTPS
    session_cookie_secure=False, # Set True in production with HTTPS
    session_cookie_http_only=True,
    frame_options='DENY',
    strict_transport_security=False, # Set True in production with HTTPS
)


# === Constants ===
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'ppt', 'pptx', 'jpg', 'jpeg'}
ALLOWED_IMAGE_EXTENSIONS = {'pdf', 'jpg', 'jpeg'}
ALLOWED_MIME_TYPES = {
    'pdf': ['application/pdf'],
    'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword'],
    'ppt': ['application/vnd.ms-powerpoint'],
    'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
    'jpg': ['image/jpeg'],
    'jpeg': ['image/jpeg']
}
LIBREOFFICE_TIMEOUT = 180
GS_TIMEOUT = 180
MIME_BUFFER_SIZE = 4096

# === Helper Functions ===
def make_error_response(error_key, status_code=400):
    logger.warning(f"Returning error: {error_key} (Status: {status_code})")
    response_text = f"Conversion failed: {error_key}"
    response = make_response(response_text, status_code)
    response.headers["Content-Type"] = "text/plain; charset=utf-8"
    return response

# --- Logic tìm và xác minh LibreOffice (ĐÃ SỬA ĐỂ KIỂM TRA PATH WINDOWS TRƯỚC) ---
SOFFICE_PATH = None
potential_paths_win = []

# Chỉ kiểm tra đường dẫn Windows chuẩn trên Windows
if sys.platform == 'win32':
    common_paths = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe'
    ]
    potential_paths_win.extend(common_paths)

    # Thử tìm trong các đường dẫn Windows chuẩn trước
    for path_to_check in potential_paths_win:
        if os.path.isfile(path_to_check):
            logger.info(f"Found potential LO path by checking known Windows location: {path_to_check}")
            # TRÊN WINDOWS: Chỉ cần file tồn tại là đủ tin tưởng, bỏ qua version check
            SOFFICE_PATH = path_to_check
            logger.info(f"Assuming valid LO on Windows (path exists) and setting path: {SOFFICE_PATH}")
            break # Đã tìm thấy, thoát vòng lặp

    # Nếu chưa tìm thấy ở path chuẩn, thử shutil.which trên Windows
    if not SOFFICE_PATH:
        logger.info("Could not find LO in known Windows locations, trying shutil.which('soffice.exe')...")
        soffice_found_which = shutil.which('soffice.exe') # Ưu tiên tìm .exe trên Windows
        if soffice_found_which:
             # TRÊN WINDOWS: Chỉ cần which tìm thấy là đủ tin tưởng
             SOFFICE_PATH = soffice_found_which
             logger.info(f"Found LO via shutil.which on Windows. Assuming valid and setting path: {SOFFICE_PATH}")
        else:
             logger.warning("shutil.which('soffice.exe') did not find an executable on Windows.")

# Logic cho các hệ điều hành khác (Linux/Mac - Giữ nguyên version check)
elif not SOFFICE_PATH: # Dùng elif để chỉ chạy phần này nếu không phải Windows VÀ chưa tìm thấy
    logger.info("Not Windows or LO not found yet, trying shutil.which('libreoffice')...")
    soffice_found_which = shutil.which('libreoffice')
    if soffice_found_which:
        try:
            # Giữ lại version check cho Non-Windows
            version_cmd = [soffice_found_which, '--headless', '--version'] # Dùng --headless ở đây
            result = subprocess.run(version_cmd, capture_output=True, text=True, check=False, timeout=15)
            if result.returncode == 0 and 'LibreOffice' in result.stdout:
                logger.info(f"Using LO path found via shutil.which (non-Windows): {soffice_found_which}")
                SOFFICE_PATH = soffice_found_which
            else:
                logger.warning(f"Found LO path {soffice_found_which} via which (non-Windows), but version check failed! Code: {result.returncode}, Output: {result.stdout.strip()}")
        except Exception as e:
            logger.warning(f"Error verifying LO path via which (non-Windows) {soffice_found_which}: {e}")
    else:
        logger.warning("shutil.which('libreoffice') did not find an executable (non-Windows).")


# Log kết quả cuối cùng
if SOFFICE_PATH:
    logger.info(f"Successfully set LO path for use: {SOFFICE_PATH}")
else:
    # Log lỗi nghiêm trọng nếu không tìm thấy bằng mọi cách
    logger.critical("LibreOffice could NOT be set/verified using any method. Conversions requiring it WILL FAIL.")
# --- Kết thúc logic LibreOffice đã sửa ---

# --- Logic tìm và xác minh Ghostscript (ĐÃ SỬA ĐỂ KIỂM TRA PATH WINDOWS TRƯỚC) ---
GS_PATH = None
potential_gs_paths = []
gs_executable_names = [] # Danh sách các tên file .exe/lệnh có thể có

# Ưu tiên kiểm tra đường dẫn Windows chuẩn
if sys.platform == 'win32':
    # Tìm thư mục gốc GS
    gs_base_dir_pf = r'C:\Program Files\gs'
    gs_base_dir_pf86 = r'C:\Program Files (x86)\gs'
    gs_version_dir_pattern_pf = os.path.join(gs_base_dir_pf, 'gs*')
    gs_version_dir_pattern_pf86 = os.path.join(gs_base_dir_pf86, 'gs*')

    possible_gs_dirs = [d for d in glob.glob(gs_version_dir_pattern_pf) if os.path.isdir(d)]
    possible_gs_dirs.extend([d for d in glob.glob(gs_version_dir_pattern_pf86) if os.path.isdir(d)])

    if not possible_gs_dirs:
        logger.warning(f"Could not find Ghostscript version directory in {gs_base_dir_pf} or {gs_base_dir_pf86}")
    else:
        # Ưu tiên phiên bản mới nhất nếu có nhiều (sắp xếp theo tên thư mục)
        possible_gs_dirs.sort(reverse=True) # Sắp xếp giảm dần, hy vọng phiên bản cao hơn ở đầu
        latest_gs_dir = possible_gs_dirs[0]
        bin_path = os.path.join(latest_gs_dir, 'bin')
        if os.path.isdir(bin_path):
            # Thêm các đường dẫn đầy đủ tiềm năng
            potential_gs_paths.append(os.path.join(bin_path, 'gswin64c.exe')) # Ưu tiên console 64bit
            potential_gs_paths.append(os.path.join(bin_path, 'gswin32c.exe')) # Console 32bit
            potential_gs_paths.append(os.path.join(bin_path, 'gs.exe'))       # Tên gốc
            # Lưu lại các tên exe để dùng với shutil.which nếu cần
            gs_executable_names.extend(['gswin64c.exe', 'gswin32c.exe', 'gs.exe', 'gs'])
            logger.info(f"Checking for Ghostscript in potential directory: {bin_path}")
        else:
            logger.warning(f"Ghostscript bin directory not found in {latest_gs_dir}")
# Đặt tên mặc định cho Linux/Mac
if not gs_executable_names:
    gs_executable_names.append('gs')

# Thử tìm trong các đường dẫn tiềm năng trước (Windows)
for path_to_check in potential_gs_paths:
    if os.path.isfile(path_to_check):
        logger.info(f"Found potential GS path by checking known location: {path_to_check}")
        try:
            result = subprocess.run([path_to_check, '--version'], capture_output=True, text=True, check=False, timeout=10)
            if result.returncode == 0 and '.' in result.stdout.strip():
                logger.info(f"Verified GS path by checking known location: {path_to_check} (Version: {result.stdout.strip()})")
                GS_PATH = path_to_check
                break # Đã tìm thấy và xác minh
            else:
                logger.warning(f"Known GS path {path_to_check} exists, but version check failed! Code: {result.returncode}, Output: {result.stdout.strip()}")
        except Exception as e:
            logger.warning(f"Error verifying known GS path {path_to_check}: {e}")
    if GS_PATH: break # Thoát sớm nếu đã tìm thấy

# Nếu vẫn chưa tìm thấy, thử dùng shutil.which với các tên khả thi
if not GS_PATH:
     logger.info(f"Could not verify GS in known locations, trying shutil.which with {gs_executable_names}...")
     for gs_name in gs_executable_names:
         gs_found_which = shutil.which(gs_name)
         if gs_found_which:
             try:
                 result = subprocess.run([gs_found_which, '--version'], capture_output=True, text=True, check=False, timeout=10)
                 if result.returncode == 0 and '.' in result.stdout.strip():
                     logger.info(f"Using GS path found via shutil.which('{gs_name}'): {gs_found_which} (Version: {result.stdout.strip()})")
                     GS_PATH = gs_found_which
                     break # Tìm thấy qua which
                 else:
                      logger.warning(f"Found GS path {gs_found_which} via which('{gs_name}'), but version check failed! Code: {result.returncode}, Output: {result.stdout.strip()}")
             except Exception as e:
                  logger.warning(f"Error verifying GS path via which('{gs_name}') {gs_found_which}: {e}")
         # else: logger.debug(f"shutil.which('{gs_name}') did not find an executable.")
     if not GS_PATH: logger.warning(f"shutil.which failed for all potential names: {gs_executable_names}")


# Log kết quả cuối cùng
if not GS_PATH:
    logger.critical("Ghostscript could NOT be found or verified. PDF Compression WILL FAIL.")
else:
    logger.info(f"Successfully set GS path for use: {GS_PATH}")
# --- Kết thúc logic Ghostscript đã sửa ---


def _allowed_file_extension(filename, allowed_set):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_set

def safe_remove(item_path, retries=3, delay=0.5):
    # Giữ nguyên hàm này
    if not item_path or not os.path.exists(item_path): return True
    is_dir = os.path.isdir(item_path); item_type = "directory" if is_dir else "file"
    for i in range(retries):
        try:
            if is_dir: shutil.rmtree(item_path)
            else: os.remove(item_path)
            logger.debug(f"Removed {item_type}: {item_path}"); return True
        except Exception as e: logger.warning(f"Error removing {item_path} (Attempt {i+1}): {e}"); time.sleep(delay*(i+1))
    logger.error(f"Failed to remove {item_type} after {retries} attempts: {item_path}"); return False

def get_actual_mime_type(file_storage):
    # Giữ nguyên hàm này, nhưng thêm kiểm tra magic đã import thành công chưa
    if not magic:
        logger.warning("python-magic not available, skipping MIME type validation.")
        # Có thể trả về None hoặc một giá trị giả định nếu cần
        # Hoặc dựa vào phần mở rộng file như một phương án dự phòng yếu hơn
        filename = file_storage.filename
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        # Giả định dựa trên phần mở rộng - KHÔNG AN TOÀN BẰNG MAGIC
        if ext == 'pdf': return 'application/pdf'
        if ext == 'docx': return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        # ... thêm các giả định khác nếu muốn, nhưng không khuyến khích
        return None # Hoặc raise lỗi

    mime_type = None
    try:
        original_pos = file_storage.stream.tell()
        file_storage.stream.seek(0)
        buffer = file_storage.stream.read(MIME_BUFFER_SIZE)
        file_storage.stream.seek(original_pos)
        mime_type = magic.from_buffer(buffer, mime=True)
        logger.debug(f"Detected MIME type: {mime_type} for file {file_storage.filename}")
    except magic.MagicException as e: logger.warning(f"Could not determine MIME type for {file_storage.filename}: {e}")
    except Exception as e: logger.error(f"Unexpected error during MIME detection for {file_storage.filename}: {e}")
    return mime_type

# --- Other Helper Functions ---
# Giữ nguyên các hàm: get_pdf_page_size, setup_slide_size, sort_key_for_pptx_images,
# _convert_pdf_to_pptx_images, convert_pdf_to_pptx_python, convert_images_to_pdf,
# convert_pdf_to_image_zip, compress_pdf_ghostscript

def get_pdf_page_size(pdf_path):
    width, height = None, None
    try:
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f);
            if reader.is_encrypted:
                try:
                    # Thử mật khẩu rỗng trước
                    decrypt_result = reader.decrypt('')
                    # Kiểm tra kết quả trả về (có thể khác nhau giữa các phiên bản)
                    if isinstance(decrypt_result, PyPDF2.PasswordType) and decrypt_result == PyPDF2.PasswordType.UNKNOWN_PASSWORD:
                        logger.warning(f"PDF is password protected (empty password failed): {pdf_path}")
                        raise ValueError("err-pdf-protected")
                    elif decrypt_result == 0: # Một số phiên bản cũ trả về 0 nếu sai pass
                         logger.warning(f"PDF is password protected (decrypt returned 0): {pdf_path}")
                         raise ValueError("err-pdf-protected")
                    elif not reader.pages: # Nếu decrypt thành công nhưng vẫn ko đọc đc trang
                        logger.warning(f"PDF possibly password protected (decrypt might have partially worked but no pages): {pdf_path}")
                        raise ValueError("err-pdf-protected")

                except ValueError as ve_decrypt: # Bắt lỗi err-pdf-protected đã raise
                    raise ve_decrypt
                except NotImplementedError:
                     logger.warning(f"Decryption algorithm not supported by PyPDF2 for {pdf_path}. Assuming protected.")
                     raise ValueError("err-pdf-protected")
                except Exception as decrypt_err:
                    logger.warning(f"Decryption attempt failed for {pdf_path}: {decrypt_err}")
                    raise ValueError("err-pdf-protected")

            if not reader.pages:
                 # Kiểm tra lại sau khi decrypt có thể thành công
                 if reader.is_encrypted: # Nếu vẫn encrypted -> lỗi decrypt
                      raise ValueError("err-pdf-protected")
                 else: # Không encrypted mà không có trang -> lỗi file
                      raise ValueError("err-pdf-no-pages")

            page = reader.pages[0]; box = page.mediabox or page.cropbox
            if box: width, height = float(box.width), float(box.height)
    except PyPDF2.errors.PdfReadError as pdf_err: raise ValueError("err-pdf-corrupt") from pdf_err
    except ValueError as ve: raise ve # Raise lại lỗi đã xác định (protected, no pages, corrupt)
    except Exception as e: logger.error(f"Error reading PDF size {pdf_path}: {e}"); raise ValueError("err-unknown") # Lỗi không xác định khác
    return width, height

def setup_slide_size(prs, pdf_path):
    try:
        pdf_width_pt, pdf_height_pt = get_pdf_page_size(pdf_path)
        if pdf_width_pt is None or pdf_height_pt is None: # Kiểm tra cả hai
            logger.warning("Could not get PDF page size for slide setup. Falling back.")
            prs.slide_width, prs.slide_height = Inches(10), Inches(7.5) # Default 4:3
        else:
            pdf_width_in, pdf_height_in = pdf_width_pt / 72.0, pdf_height_pt / 72.0; max_dim = 56.0 # Max PPTX slide dim in inches
            # Scale down if too large, preserving aspect ratio
            if pdf_width_in > max_dim or pdf_height_in > max_dim:
                ratio = pdf_width_in / pdf_height_in if pdf_height_in > 0 else 1
                if pdf_width_in >= pdf_height_in:
                    final_width = max_dim
                    final_height = max_dim / ratio if ratio != 0 else max_dim # Avoid division by zero
                else:
                    final_height = max_dim
                    final_width = max_dim * ratio
                logger.info(f"PDF dims ({pdf_width_in:.2f}x{pdf_height_in:.2f}) too large, scaled to {final_width:.2f}x{final_height:.2f}")
            else:
                final_width, final_height = pdf_width_in, pdf_height_in
            # Ensure dimensions are positive
            if final_width <= 0 or final_height <= 0:
                 logger.warning(f"Calculated non-positive slide dimensions ({final_width}x{final_height}). Falling back.")
                 prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
            else:
                 prs.slide_width, prs.slide_height = Inches(final_width), Inches(final_height)
                 logger.info(f"Set slide size from PDF: {final_width:.2f}in x {final_height:.2f}in")

    except ValueError as ve:
         # Nếu get_pdf_page_size báo lỗi cụ thể, log và fallback
         logger.warning(f"Error getting PDF page size ({ve}). Falling back on slide setup.")
         prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    except Exception as e:
        logger.warning(f"Error setting slide size from PDF dims: {e}. Falling back.")
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    return prs

def sort_key_for_pptx_images(filename):
    # Giữ nguyên
    try: return int(os.path.splitext(filename)[0].split('-')[-1].split('_')[-1])
    except (ValueError, IndexError): return 0

def _convert_pdf_to_pptx_images(input_path, output_path):
    # Giữ nguyên logic cốt lõi, chỉ đảm bảo lỗi được raise rõ ràng
    temp_dir = None
    success = False
    try:
        temp_dir = tempfile.mkdtemp(prefix="pdfimg_")
        page_count = 0 # Khởi tạo
        try:
            # Cố gắng lấy page count và kiểm tra mã hóa
            page_count_info = pdfinfo_from_path(input_path, poppler_path=None)
            page_count = page_count_info.get('Pages')
            # Nếu poppler không lấy được page count, thử với PyPDF2 như một backup
            if page_count is None:
                logger.warning("pdfinfo failed to get page count, trying PyPDF2...")
                try:
                    with open(input_path, 'rb') as f:
                         reader = PyPDF2.PdfReader(f, strict=False) # strict=False có thể bỏ qua lỗi nhỏ
                         if reader.is_encrypted:
                             # Thử decrypt rỗng, nếu lỗi -> protected
                             try:
                                 reader.decrypt('')
                                 if not reader.pages: # Vẫn ko đọc được -> protected
                                     raise ValueError("err-pdf-protected-pypdf2")
                             except Exception:
                                  raise ValueError("err-pdf-protected-pypdf2")
                         page_count = len(reader.pages)
                    logger.info(f"PyPDF2 got page count: {page_count}")
                except ValueError as ve_pypdf:
                     if "err-pdf-protected" in str(ve_pypdf): raise ve_pypdf # Re-raise lỗi protected
                     else: raise PDFPageCountError(f"PyPDF2 failed to get page count: {ve_pypdf}") from ve_pypdf
                except Exception as e_pypdf:
                     raise PDFPageCountError(f"PyPDF2 failed to get page count: {e_pypdf}") from e_pypdf

            # Nếu cả hai cách đều ko lấy được page count -> lỗi
            if page_count is None:
                 raise PDFPageCountError("Could not determine page count using pdfinfo or PyPDF2.")

            # Kiểm tra mã hóa bằng pdfinfo nếu có thể, hoặc dựa vào get_pdf_page_size
            is_encrypted = page_count_info.get('Encrypted', 'no').lower() == 'yes'
            if is_encrypted:
                 try:
                     # Dùng hàm đã có để kiểm tra decrypt rỗng
                     get_pdf_page_size(input_path)
                 except ValueError as ve_enc:
                      if "err-pdf-protected" in str(ve_enc):
                          raise ve_enc # Raise lại lỗi protected đã xác định
                      else: # Lỗi khác khi đọc size PDF mã hóa?
                           logger.warning(f"Error checking encrypted PDF with get_pdf_page_size: {ve_enc}")
                           raise # Raise lỗi đó lên
                 except Exception as e_enc:
                      logger.warning(f"Unexpected error checking encrypted PDF: {e_enc}")
                      raise # Raise lỗi đó lên

        # Xử lý các lỗi cụ thể từ pdfinfo/PyPDF2
        except PDFInfoNotInstalledError as e: raise ValueError("err-poppler-missing") from e
        except (PDFPageCountError, PDFSyntaxError) as e: raise ValueError("err-pdf-corrupt") from e
        except ValueError as ve: # Bắt lỗi err-pdf-protected từ các block try bên trong
             if "err-pdf-protected" in str(ve): raise ve
             else: raise # Raise lỗi ValueError khác nếu có
        except Exception as info_err: # Lỗi không mong muốn khác khi lấy info
             logger.error(f"Unexpected error getting PDF info for PPTX conversion: {info_err}", exc_info=True)
             raise RuntimeError("err-conversion") from info_err


        if page_count == 0:
             logger.info("PDF has 0 pages. Creating empty PPTX.")
             Presentation().save(output_path)
             success = True
        else:
            logger.info(f"Converting {page_count} PDF pages to images for PPTX...")
            # Sử dụng thread_count = 1 có thể ổn định hơn trên một số hệ thống
            images = convert_from_path(input_path, dpi=300, fmt='jpeg', output_folder=temp_dir, thread_count=1, poppler_path=None, strict=False)
            if not images:
                 # Kiểm tra lại page_count phòng trường hợp convert_from_path lỗi âm thầm
                 if page_count > 0:
                     logger.error("convert_from_path returned no images despite page count > 0.")
                     raise RuntimeError("err-conversion-img")
                 else: # Trường hợp hy hữu page_count=0 nhưng images=None
                      logger.warning("No images generated for 0-page PDF.")
                      Presentation().save(output_path)
                      success = True

            if images: # Chỉ xử lý nếu có images trả về
                prs = Presentation()
                prs = setup_slide_size(prs, input_path) # Thiết lập kích thước slide
                blank_layout = prs.slide_layouts[6] # Layout trống

                # Sắp xếp ảnh đã tạo (pdf2image thường tạo tên file có dạng ...-01.jpg, ...)
                gen_imgs = sorted([f for f in os.listdir(temp_dir) if f.lower().endswith(('.jpg', '.jpeg'))], key=sort_key_for_pptx_images)

                if not gen_imgs:
                    logger.error("pdf2image ran but no image files found in temp directory.")
                    raise RuntimeError("err-conversion-img")

                logger.info(f"Adding {len(gen_imgs)} generated images to PPTX...")
                slide_w, slide_h = prs.slide_width, prs.slide_height

                for img_fn in gen_imgs:
                    img_path = os.path.join(temp_dir, img_fn)
                    try:
                        slide = prs.slides.add_slide(blank_layout)
                        # Lấy kích thước ảnh và tính toán vị trí/kích thước để vừa slide
                        with Image.open(img_path) as img:
                            img_w_px, img_h_px = img.size
                        if img_w_px <= 0 or img_h_px <= 0:
                             logger.warning(f"Image {img_fn} has zero dimension. Skipping.")
                             continue

                        img_aspect_ratio = img_w_px / img_h_px
                        slide_aspect_ratio = slide_w / slide_h if slide_h > 0 else 1

                        # Tính toán kích thước ảnh trên slide để giữ tỷ lệ và vừa khung
                        if img_aspect_ratio > slide_aspect_ratio:
                            # Ảnh rộng hơn slide -> chiều rộng ảnh = chiều rộng slide
                            pic_w = slide_w
                            pic_h = int(slide_w / img_aspect_ratio)
                        else:
                            # Ảnh cao hơn hoặc bằng slide -> chiều cao ảnh = chiều cao slide
                            pic_h = slide_h
                            pic_w = int(slide_h * img_aspect_ratio)

                        # Đảm bảo kích thước > 0
                        if pic_w <= 0 or pic_h <= 0:
                             logger.warning(f"Calculated zero dimension for {img_fn} on slide. Skipping.")
                             continue

                        # Căn giữa ảnh
                        pic_l = int((slide_w - pic_w) / 2)
                        pic_t = int((slide_h - pic_h) / 2)

                        slide.shapes.add_picture(img_path, pic_l, pic_t, width=pic_w, height=pic_h)

                    except UnidentifiedImageError:
                        logger.warning(f"Skipping invalid image file {img_fn}")
                        continue
                    except Exception as page_err:
                        logger.warning(f"Error adding image {img_fn} to PPTX slide: {page_err}")
                        # Có thể quyết định dừng lại hoặc tiếp tục với các ảnh khác
                        # continue

                prs.save(output_path)
                logger.info("PPTX file created successfully.")
                success = True

    except ValueError as ve: logger.error(f"PDF->PPTX(Image) Value Error: {ve}"); raise ve
    except RuntimeError as rte: logger.error(f"PDF->PPTX(Image) Runtime Error: {rte}"); raise rte
    except Exception as e: logger.error(f"Unexpected PDF->PPTX(Image) Error: {e}", exc_info=True); raise RuntimeError("err-unknown") from e
    finally: safe_remove(temp_dir)
    return success

def convert_pdf_to_pptx_python(input_path, output_path):
    # Giữ nguyên
    logger.info("Attempting PDF -> PPTX via Python (image-based)...")
    return _convert_pdf_to_pptx_images(input_path, output_path)

def convert_images_to_pdf(image_paths, output_path):
    # Giữ nguyên hàm này
    image_objects = []
    success = False
    try:
        sorted_paths = image_paths # Giả sử đã được sắp xếp trước nếu cần
        for file_path in sorted_paths:
            filename = os.path.basename(file_path)
            try:
                with Image.open(file_path) as img:
                    img.load() # Load image data while file is open
                    converted_img = None
                    # Handle transparency for common modes
                    if img.mode == 'RGBA':
                        logger.debug(f"Converting RGBA image {filename} to RGB with white background.")
                        bg = Image.new('RGB', img.size, (255, 255, 255))
                        try:
                            bg.paste(img, mask=img.getchannel('A'))
                            converted_img = bg
                        except Exception as paste_err:
                            logger.warning(f"Error pasting RGBA {filename}, falling back to basic convert: {paste_err}")
                            converted_img = img.convert('RGB')
                    elif img.mode == 'LA': # Luminance + Alpha
                         logger.debug(f"Converting LA image {filename} to RGB with white background.")
                         bg = Image.new('RGB', img.size, (255, 255, 255));
                         l_channel = img.getchannel('L'); a_channel = img.getchannel('A');
                         rgb_img = Image.merge('RGB', (l_channel, l_channel, l_channel))
                         try:
                             bg.paste(rgb_img, mask=a_channel); converted_img = bg
                         except Exception as paste_err:
                             logger.warning(f"Error pasting LA {filename}, falling back to basic convert: {paste_err}");
                             converted_img = img.convert('RGB')
                    elif img.mode == 'P' and 'transparency' in img.info:
                         logger.debug(f"Converting Palette image {filename} with transparency to RGB.")
                         img_rgba = img.convert('RGBA') # Convert palette with transparency to RGBA first
                         bg = Image.new('RGB', img_rgba.size, (255, 255, 255))
                         try:
                             bg.paste(img_rgba, mask=img_rgba.getchannel('A')); converted_img = bg
                         except Exception as paste_err:
                             logger.warning(f"Error pasting P->RGBA {filename}, falling back to basic convert: {paste_err}")
                             converted_img = img.convert('RGB')
                    elif img.mode not in ['RGB', 'L', 'CMYK']:
                        logger.debug(f"Converting {filename} from mode {img.mode} to RGB")
                        converted_img = img.convert('RGB')
                    else:
                        # Keep RGB, L (grayscale), CMYK as they are supported by Pillow's PDF writer
                        if img.mode == 'CMYK': logger.debug(f"Image {filename} is CMYK.")
                        if img.mode == 'L': logger.debug(f"Image {filename} is Grayscale.")
                        converted_img = img.copy() # Make a copy to avoid issues with context manager

                    if converted_img: image_objects.append(converted_img)
                    else: logger.error(f"Failed to prepare image {filename} for PDF conversion."); raise RuntimeError("err-conversion")

            except UnidentifiedImageError:
                logger.error(f"File {filename} is not a valid image or format not supported by Pillow.")
                raise ValueError("err-invalid-image-file")
            except Exception as img_err:
                logger.error(f"Error processing image {filename}: {img_err}", exc_info=True)
                raise RuntimeError("err-conversion") from img_err

        if not image_objects:
             logger.warning("No valid images found to convert to PDF.")
             raise ValueError("err-select-file") # Hoặc lỗi khác phù hợp

        # Lưu ảnh đầu tiên và nối các ảnh còn lại
        first_image = image_objects[0]
        other_images = image_objects[1:]
        first_image.save(output_path, "PDF", resolution=100.0, save_all=True, append_images=other_images)
        logger.info(f"Saved {len(image_objects)} images to PDF: {output_path}")
        success = True

    except ValueError as ve: raise ve
    except RuntimeError as rte: raise rte
    except Exception as e:
        logger.error(f"Unexpected error converting images to PDF: {e}", exc_info=True)
        raise RuntimeError("err-unknown") from e
    finally:
        # Đóng tất cả các đối tượng Image đã mở
        for img_obj in image_objects:
             try: img_obj.close()
             except Exception as close_err: logger.debug(f"Error closing PIL object: {close_err}")
    return success

def convert_pdf_to_image_zip(input_path, output_zip_path, img_format='jpeg'):
    # Giữ nguyên hàm này
    temp_dir = None; fmt = img_format.lower(); ext = 'jpg' if fmt in ['jpeg', 'jpg'] else fmt
    success = False
    try:
        temp_dir = tempfile.mkdtemp(prefix="pdf2imgzip_")
        page_count = 0 # Khởi tạo
        try:
             # Cố gắng lấy page count và kiểm tra mã hóa
             page_count_info = pdfinfo_from_path(input_path, poppler_path=None)
             page_count = page_count_info.get('Pages')
             if page_count is None:
                  logger.warning("pdfinfo failed to get page count for zip conversion, trying PyPDF2...")
                  try:
                      with open(input_path, 'rb') as f:
                           reader = PyPDF2.PdfReader(f, strict=False)
                           if reader.is_encrypted:
                               try: reader.decrypt('')
                               except Exception: raise ValueError("err-pdf-protected-pypdf2")
                               if not reader.pages: raise ValueError("err-pdf-protected-pypdf2")
                           page_count = len(reader.pages)
                      logger.info(f"PyPDF2 got page count for zip: {page_count}")
                  except ValueError as ve_pypdf:
                       if "err-pdf-protected" in str(ve_pypdf): raise ve_pypdf
                       else: raise PDFPageCountError(f"PyPDF2 failed to get page count for zip: {ve_pypdf}") from ve_pypdf
                  except Exception as e_pypdf:
                       raise PDFPageCountError(f"PyPDF2 failed to get page count for zip: {e_pypdf}") from e_pypdf

             if page_count is None: raise PDFPageCountError("Could not determine page count for zip conversion.")

             is_encrypted = page_count_info.get('Encrypted', 'no').lower() == 'yes'
             if is_encrypted:
                  try: get_pdf_page_size(input_path)
                  except ValueError as ve_enc:
                       if "err-pdf-protected" in str(ve_enc): raise ve_enc
                       else: raise
                  except Exception as e_enc: raise

             logger.info(f"PDF Info for ZIP conversion: {page_count} pages.")

        except (PDFInfoNotInstalledError, FileNotFoundError) as e: raise ValueError("err-poppler-missing") from e
        except (PDFPageCountError, PDFSyntaxError) as e: raise ValueError("err-pdf-corrupt") from e
        except ValueError as ve: # Bắt lỗi err-pdf-protected
             if "err-pdf-protected" in str(ve): raise ve
             else: raise
        except Exception as info_err: logger.error(f"pdfinfo error for zip conversion: {info_err}"); raise ValueError("err-poppler-check-failed") from info_err

        if page_count == 0:
            logger.warning("PDF has 0 pages. Creating empty ZIP.")
            with zipfile.ZipFile(output_zip_path, 'w') as zf: pass
            success = True
        else:
            # Tạo base name an toàn hơn
            input_basename = os.path.splitext(os.path.basename(input_path))[0]
            safe_base = secure_filename(f"page_{input_basename}")[:100] # Giới hạn độ dài

            images = convert_from_path(input_path, dpi=200, fmt=fmt, output_folder=temp_dir, output_file=safe_base, thread_count=1, poppler_path=None, strict=False)
            if not images:
                 if page_count > 0: logger.error(f"pdf2image failed for zip conversion {input_path}"); raise RuntimeError("err-conversion-img")
                 else:
                     logger.warning("PDF 0 pages, no images generated. Empty ZIP.")
                     with zipfile.ZipFile(output_zip_path, 'w') as zf: pass
                     success = True
            else:
                def sort_key_pdf2image(f):
                    # Cải thiện key sort để xử lý tên file phức tạp hơn
                    try:
                        # Tìm số cuối cùng trong tên file (sau dấu gạch ngang hoặc gạch dưới)
                        parts = os.path.splitext(f)[0].replace('-', '_').split('_')
                        for part in reversed(parts):
                            if part.isdigit():
                                return int(part)
                        return 0 # Không tìm thấy số
                    except (ValueError, IndexError):
                        logger.warning(f"Cannot get page number from '{f}' for sorting")
                        return 0

                # Lọc file chính xác hơn dựa trên prefix đã tạo
                gen_files = sorted( [f for f in os.listdir(temp_dir) if f.lower().startswith(safe_base.lower()) and f.lower().endswith(f'.{ext}')], key=sort_key_pdf2image )

                if not gen_files and page_count > 0:
                     logger.error(f"No output files found matching pattern '{safe_base}*.{ext}' in {temp_dir}")
                     # Thử list lại xem có file nào không
                     all_files = os.listdir(temp_dir)
                     logger.debug(f"Files found in temp dir: {all_files}")
                     raise RuntimeError("err-conversion-img")
                elif not gen_files and page_count == 0:
                     logger.warning("PDF 0 pages, no images found after conversion attempt. Empty ZIP.")
                     with zipfile.ZipFile(output_zip_path, 'w') as zf: pass
                     success = True
                else:
                    with zipfile.ZipFile(output_zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                         for i, filename in enumerate(gen_files):
                             zip_filename = f"page_{i+1}.{ext}" # Đặt tên file trong zip đơn giản
                             zf.write(os.path.join(temp_dir, filename), zip_filename)
                    logger.info(f"Created image ZIP: {output_zip_path} with {len(gen_files)} images.")
                    success = True
    except ValueError as ve: raise ve
    except RuntimeError as rte: raise rte
    except Exception as e: logger.error(f"Unexpected PDF->ZIP Error: {e}", exc_info=True); raise RuntimeError("err-unknown") from e
    finally: safe_remove(temp_dir)
    return success

def compress_pdf_ghostscript(input_path, output_path, quality_level='medium'):
    # Giữ nguyên hàm này
    if not GS_PATH: logger.error("GS_PATH not set."); raise RuntimeError("err-gs-missing")
    success = False
    # Thêm check file input tồn tại
    if not os.path.isfile(input_path):
         logger.error(f"Input file for GS compression not found: {input_path}")
         raise RuntimeError("err-gs-failed") # Hoặc lỗi khác

    # Sử dụng GS_PATH đã được xác minh
    gs_base_cmd = [ GS_PATH, '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dNOPAUSE', '-dBATCH', '-dQUIET' ]
    gs_output_cmd = [f'-sOutputFile={output_path}']
    gs_input_cmd = [input_path]
    cmd = []; log_quality_info = ""

    # Chọn cài đặt chất lượng
    if quality_level == 'low':
        ppi = 120 # Hoặc 72 tùy mức độ nén mong muốn
        specific_settings = [
            '-dDownsampleColorImages=true', '-dDownsampleGrayImages=true', '-dDownsampleMonoImages=true',
            f'-dColorImageResolution={ppi}', f'-dGrayImageResolution={ppi}', f'-dMonoImageResolution={ppi}',
            '-dColorImageDownsampleType=/Bicubic', '-dGrayImageDownsampleType=/Bicubic', '-dMonoImageDownsampleType=/Bicubic',
            '-dEmbedAllFonts=true', # Giữ lại để tương thích tốt hơn
            '-dSubsetFonts=true',
            '-dAutoRotatePages=/None' # Tránh xoay trang không mong muốn
            #'-dDetectDuplicateImages=true' # Có thể thêm để giảm size nếu có ảnh trùng lặp
        ]
        cmd = gs_base_cmd + specific_settings + gs_output_cmd + gs_input_cmd
        log_quality_info = f"Quality: low (Target PPI: {ppi})"
    elif quality_level == 'high':
         # Cài đặt '/printer' chất lượng cao nhất
         standard_settings = [f'-dPDFSETTINGS=/printer']
         cmd = gs_base_cmd + standard_settings + gs_output_cmd + gs_input_cmd
         log_quality_info = f"Quality: high (Setting: /printer)"
    else: # Mặc định là medium ('/ebook')
        standard_settings = [f'-dPDFSETTINGS=/ebook']
        cmd = gs_base_cmd + standard_settings + gs_output_cmd + gs_input_cmd
        log_quality_info = f"Quality: medium (Setting: /ebook)"

    logger.info(f"Running Ghostscript ({log_quality_info}): {' '.join(cmd)}")
    try:
        # Chạy subprocess
        result = subprocess.run(cmd, check=True, timeout=GS_TIMEOUT, capture_output=True, text=True, encoding='utf-8', errors='ignore')
        if result.stdout: logger.info(f"Ghostscript stdout:\n{result.stdout}")
        if result.stderr: logger.info(f"Ghostscript stderr:\n{result.stderr}") # GS thường ghi thông tin vào stderr

        # Kiểm tra file output
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            original_size = os.path.getsize(input_path)
            compressed_size = os.path.getsize(output_path)
            # Cảnh báo nếu file output không nhỏ hơn đáng kể hoặc lớn hơn file gốc
            if compressed_size >= original_size * 0.98 :
                 logger.warning(f"GS produced minimal size reduction or file increase. Original: {original_size}, Compressed: {compressed_size}. This might happen with already optimized PDFs.")
                 # Quyết định xem có coi đây là thành công hay không. Thường thì vẫn là thành công.
            logger.info(f"GS compression successful: {output_path} (Size: {compressed_size} bytes)")
            success = True
        else:
            logger.error(f"Ghostscript command ran but output file '{output_path}' is missing or empty.")
            raise RuntimeError("err-gs-failed")

    except subprocess.TimeoutExpired:
        logger.error(f"Ghostscript command timed out ({GS_TIMEOUT}s) for {input_path}.")
        safe_remove(output_path) # Xóa file output nếu có (thường là rỗng hoặc chưa hoàn thành)
        raise RuntimeError("err-gs-timeout")
    except subprocess.CalledProcessError as gs_err:
        logger.error(f"Ghostscript command failed for {input_path}. Return Code: {gs_err.returncode}")
        # Ghi lại stdout/stderr để debug
        if gs_err.stdout: logger.error(f"GS stdout on error:\n{gs_err.stdout}")
        if gs_err.stderr: logger.error(f"GS stderr on error:\n{gs_err.stderr}")
        safe_remove(output_path) # Xóa file output lỗi nếu có

        # Phân tích stderr để xác định lỗi cụ thể hơn
        error_output = (gs_err.stderr or "").lower()
        if "password required" in error_output or "owner password required" in error_output:
            logger.warning(f"Ghostscript failed due to encrypted PDF: {input_path}")
            raise ValueError("err-pdf-protected")
        elif "corrupt" in error_output or "invalid" in error_output or "error" in error_output:
             logger.warning(f"Ghostscript indicates corrupt/invalid PDF or processing error: {input_path}")
             # Phân loại rõ hơn nếu có thể
             if "undefined" in error_output or "syntaxerror" in error_output:
                  raise ValueError("err-pdf-corrupt")
             else: # Lỗi xử lý chung của GS
                  raise RuntimeError("err-gs-failed")
        else: # Lỗi không xác định từ GS
             raise RuntimeError("err-gs-failed")

    except ValueError as ve: raise ve # Re-raise lỗi đã xác định (protected, corrupt)
    except FileNotFoundError:
        logger.error(f"Ghostscript executable not found at the specified path: {GS_PATH}")
        raise RuntimeError("err-gs-missing")
    except Exception as gs_run_err:
        logger.error(f"Unexpected error running Ghostscript: {gs_run_err}", exc_info=True)
        safe_remove(output_path)
        raise RuntimeError("err-gs-failed") # Lỗi chung khi chạy GS

    return success


# === Global Error Handlers ===
# Giữ nguyên
@app.errorhandler(CSRFError)
def handle_csrf_error(e): logger.warning(f"CSRF failed: {e.description}"); return make_error_response("err-csrf-invalid", 400)
@app.errorhandler(RequestEntityTooLarge)
def handle_file_too_large(e): logger.warning(f"File too large: {e.description}"); return make_error_response("err-file-too-large", 413)
@app.errorhandler(429)
def ratelimit_handler(e): logger.warning(f"Rate limit exceeded: {e.description}"); return make_error_response("err-rate-limit-exceeded", 429)
@app.errorhandler(Exception)
def handle_generic_exception(e):
     from werkzeug.exceptions import HTTPException
     # Nếu là lỗi HTTP đã biết (như 404, 405), để Flask xử lý mặc định hoặc trả về e
     if isinstance(e, HTTPException):
         # Log lỗi cụ thể nếu muốn
         logger.warning(f"Caught HTTPException: {e.code} - {e.name} - {e.description}")
         # Có thể tùy chỉnh response ở đây nếu muốn, hoặc return e để dùng handler mặc định
         # return make_error_response(f"err-http-{e.code}", e.code)
         return e # Trả về đối tượng lỗi HTTP để Flask xử lý tiếp
     # Nếu là lỗi không mong muốn khác
     logger.error(f"Unhandled Exception: {e}", exc_info=True)
     return make_error_response("err-unknown", 500)


# === Routes ===

# Giữ nguyên các route: /api/translations, /, /convert, /convert_image, /compress_pdf, /compress_docx
# Các thay đổi chính đã được thực hiện trong logic tìm kiếm LO/GS ở đầu file.
# Logic bên trong các route này sử dụng các biến SOFFICE_PATH và GS_PATH đã được xác định.

@app.route('/api/translations')
def get_translations():
    # Giữ nguyên
    translations = {
        'en': {
            'lang-title': 'PDF & Office Tools', 'lang-subtitle': 'Simple, powerful tools for your documents',
            'lang-error-title': 'Error!', 'lang-convert-title': 'Convert PDF/Office',
            'lang-convert-desc': 'Transform PDF to Word/PPT and vice versa',
            'lang-compress-title': 'Compress PDF', 'lang-compress-desc': 'Reduce PDF file size while optimizing for quality',
            'lang-compress-input-label': 'Select PDF file', 'lang-compress-btn': 'Compress PDF',
            'lang-compressing': 'Compressing PDF...', 'lang-select-quality': 'Compression Level',
            'lang-quality-low': 'Low Quality (Smallest Size)',
            'lang-quality-medium': 'Medium Quality (Good Balance)',
            'lang-quality-high': 'High Quality (Less Compression)',
            'lang-merge-title': 'Merge PDF', 'lang-merge-desc': 'Combine multiple PDFs into one file',
            'lang-split-title': 'Split PDF', 'lang-split-desc': 'Extract pages from your PDF',
            'lang-image-title': 'PDF ↔ Image', 'lang-image-desc': 'Convert PDF to images or images to PDF',
            'lang-image-input-label': 'Select PDF or Image(s) (JPG/JPEG only)', 'lang-image-convert-btn': 'Convert Now',
            'lang-image-converting': 'Converting...', 'lang-size-limit': 'Size limit: 100MB',
            'lang-size-limit-total': 'Size limit: 100MB (total)', 'lang-select-conversion': 'Select conversion type',
            'lang-converting': 'Converting...', 'lang-convert-btn': 'Convert Now',
            'lang-file-input-label': 'Select file', 'file-no-selected': 'No file selected',
            'err-select-file': 'Please select file(s).', 'err-file-too-large': 'File size exceeds the limit (100MB).',
            'err-select-conversion': 'Please select a conversion type.',
            'err-format-docx': 'Select one DOCX file for this operation.',
            'err-format-ppt': 'Select one PDF, PPT or PPTX file for this conversion.',
            'err-format-pdf': 'Please select a PDF file.', 'err-conversion': 'An error occurred during processing.',
            'err-fetch-translations': 'Could not load language data.', 'lang-select-btn-text': 'Browse',
            'lang-select-conversion-label': 'Conversion Type', 'err-multi-file-not-supported': 'Multi-file selection is only supported for Image to PDF conversion.',
            'err-invalid-image-file': 'One or more selected files are not valid images (Pillow error).',
            'err-image-format': 'Invalid file type. Select PDF, JPG, or JPEG based on conversion.',
            'err-image-single-pdf': 'Please select only one PDF file to convert to images.',
            'err-image-all-images': 'If selecting multiple files, all must be JPG or JPEG to convert to PDF.',
            'err-libreoffice': 'Conversion failed (Processing engine error - LO).', 'err-conversion-timeout': 'Processing timed out.',
            'err-poppler-missing': 'PDF processing library (Poppler) missing or failed.',
            'err-pdf-corrupt': 'Could not process PDF (corrupt file?).', 'err-unknown': 'An unexpected error occurred. Please try again later.',
            'err-csrf-invalid': 'Security validation failed. Please refresh the page and try again.',
            'err-rate-limit-exceeded': 'Too many requests. Please wait a moment and try again.',
            'err-invalid-mime-type': 'Invalid file type detected. The file content does not match the expected format.',
            'err-mime-unidentified-office': "Could not identify file type, it might be non-standard. Please open your file in an Office application, press 'Save' or 'Save as' to save again and upload again.",
            'err-invalid-mime-type-image': 'Invalid image type detected. Only JPEG files are allowed for Image-to-PDF.',
            'err-pdf-protected': 'Cannot process password-protected PDF.',
            'err-poppler-check-failed': 'Failed to get PDF info (Poppler check).',
            'err-conversion-img': 'Failed to convert/extract images from PDF.',
            'err-gs-missing': 'Compression engine (Ghostscript) not available.',
            'err-gs-failed': 'Compression failed (Ghostscript error). Check if PDF is valid/not protected.',
            'err-gs-timeout': 'Compression timed out.', 'err-invalid-quality': 'Invalid compression quality selected.',
            'lang-clear-all': 'Clear All', 'lang-upload-a-file': 'Upload files',
            'lang-drag-drop': 'or drag and drop', 'lang-image-types': 'PDF, JPG, JPEG up to 100MB total',
            'lang-compress-docx-title': 'Compress Word',
            'lang-compress-docx-desc': 'Reduce Word file size while optimizing for quality',
            'lang-compress-docx-input-label': 'Select Word file',
            'lang-compressing-docx': 'Compressing Word...',
            'lang-compress-docx-btn': 'Compress Word',
        },
        'vi': {
            'lang-title': 'Công Cụ PDF & Office', 'lang-subtitle': 'Công cụ đơn giản, mạnh mẽ cho tài liệu của bạn',
            'lang-error-title': 'Lỗi!', 'lang-convert-title': 'Chuyển đổi PDF/Office',
            'lang-convert-desc': 'Chuyển đổi PDF sang Word/PPT và ngược lại',
            'lang-compress-title': 'Nén PDF', 'lang-compress-desc': 'Giảm dung lượng tệp PDF mà vẫn tối ưu chất lượng',
            'lang-compress-input-label': 'Chọn tệp PDF', 'lang-compress-btn': 'Nén PDF',
            'lang-compressing': 'Đang nén PDF...', 'lang-select-quality': 'Mức độ nén',
            'lang-quality-low': 'Nén Mạnh (Nhẹ Nhất)',
            'lang-quality-medium': 'Nén Vừa (Cân Bằng)',
            'lang-quality-high': 'Nén Nhẹ (Nén Ít)',
            'lang-merge-title': 'Gộp PDF', 'lang-merge-desc': 'Kết hợp nhiều tệp PDF thành một tệp',
            'lang-split-title': 'Tách PDF', 'lang-split-desc': 'Trích xuất các trang từ tệp PDF của bạn',
            'lang-image-title': 'PDF ↔ Ảnh', 'lang-image-desc': 'Chuyển PDF thành ảnh hoặc ảnh thành PDF',
            'lang-image-input-label': 'Chọn PDF hoặc (các) Ảnh (chỉ JPG/JPEG)', 'lang-image-convert-btn': 'Chuyển đổi ngay',
            'lang-image-converting': 'Đang chuyển đổi...', 'lang-size-limit': 'Giới hạn kích thước: 100MB',
            'lang-size-limit-total': 'Giới hạn kích thước: 100MB (tổng)', 'lang-select-conversion': 'Chọn kiểu chuyển đổi',
            'lang-converting': 'Đang chuyển đổi...', 'lang-convert-btn': 'Chuyển đổi ngay',
            'lang-file-input-label': 'Chọn tệp', 'file-no-selected': 'Không có tệp nào được chọn',
            'err-select-file': 'Vui lòng chọn (các) tệp.', 'err-file-too-large': 'Kích thước tệp vượt quá giới hạn (100MB).',
            'err-select-conversion': 'Vui lòng chọn kiểu chuyển đổi.',
            'err-format-docx': 'Chọn một file DOCX cho thao tác này.',
            'err-format-ppt': 'Chọn một file PDF, PPT hoặc PPTX cho chuyển đổi này.',
            'err-format-pdf': 'Vui lòng chọn một tệp PDF.', 'err-conversion': 'Đã xảy ra lỗi trong quá trình xử lý.',
            'err-fetch-translations': 'Không thể tải dữ liệu ngôn ngữ.', 'lang-select-btn-text': 'Duyệt...',
            'lang-select-conversion-label': 'Kiểu chuyển đổi', 'err-multi-file-not-supported': 'Chỉ hỗ trợ chọn nhiều file khi chuyển đổi Ảnh sang PDF.',
            'err-invalid-image-file': 'Một hoặc nhiều tệp được chọn không phải là ảnh hợp lệ (lỗi Pillow).',
            'err-image-format': 'Loại tệp không hợp lệ. Chọn PDF, JPG, hoặc JPEG tùy theo chuyển đổi.',
            'err-image-single-pdf': 'Vui lòng chỉ chọn một file PDF để chuyển đổi sang ảnh.',
            'err-image-all-images': 'Nếu chọn nhiều tệp, tất cả phải là JPG hoặc JPEG để chuyển đổi sang PDF.',
            'err-libreoffice': 'Chuyển đổi thất bại (Lỗi bộ xử lý - LO).', 'err-conversion-timeout': 'Quá trình xử lý quá thời gian.',
            'err-poppler-missing': 'Thiếu hoặc lỗi thư viện xử lý PDF (Poppler).',
            'err-pdf-corrupt': 'Không thể xử lý PDF (tệp lỗi?).', 'err-unknown': 'Đã xảy ra lỗi không mong muốn. Vui lòng thử lại sau.',
            'err-csrf-invalid': 'Xác thực bảo mật thất bại. Vui lòng tải lại trang và thử lại.',
            'err-rate-limit-exceeded': 'Quá nhiều yêu cầu. Vui lòng đợi một lát và thử lại.',
            'err-invalid-mime-type': 'Phát hiện loại tệp không hợp lệ. Nội dung tệp không khớp định dạng mong đợi.',
            'err-mime-unidentified-office': "Không thể nhận dạng loại file dù có đuôi Office. Vui lòng mở file của bạn lên bằng ứng dụng Office, ấn 'Lưu' hoặc 'Lưu thành' để lưu lại bản mới và tải lên lại.",
            'err-invalid-mime-type-image': 'Phát hiện loại ảnh không hợp lệ. Chỉ cho phép tệp JPEG để chuyển đổi Ảnh sang PDF.',
            'err-pdf-protected': 'Không thể xử lý PDF được bảo vệ bằng mật khẩu.',
            'err-poppler-check-failed': 'Không thể lấy thông tin PDF (lỗi kiểm tra Poppler).',
            'err-conversion-img': 'Không thể chuyển đổi/trích xuất ảnh từ PDF.',
            'err-gs-missing': 'Không tìm thấy công cụ nén (Ghostscript).',
            'err-gs-failed': 'Nén thất bại (Lỗi Ghostscript). Kiểm tra PDF hợp lệ/không bị khóa.',
            'err-gs-timeout': 'Nén quá thời gian.', 'err-invalid-quality': 'Đã chọn mức nén không hợp lệ.',
            'lang-clear-all': 'Xóa tất cả', 'lang-upload-a-file': 'Tải tệp lên',
            'lang-drag-drop': 'hoặc kéo và thả', 'lang-image-types': 'PDF, JPG, JPEG tối đa 100MB tổng',
            'lang-compress-docx-title': 'Nén Word',
            'lang-compress-docx-desc': 'Giảm dung lượng tệp Word mà vẫn tối ưu chất lượng',
            'lang-compress-docx-input-label': 'Chọn tệp DOCX',
            'lang-compressing-docx': 'Đang nén Word...',
            'lang-compress-docx-btn': 'Nén Word',
        }
    }
    lang = request.args.get('lang', 'en')
    return jsonify(translations.get(lang, translations.get('en', {})))

@app.route('/')
def index():
    # Giữ nguyên
    try:
        translations_url = url_for('get_translations', _external=False)
        # Các biến này giờ sẽ phản ánh kết quả tìm kiếm mới
        gs_available = GS_PATH is not None
        soffice_available = SOFFICE_PATH is not None
        return render_template('index.html',
                               translations_url=translations_url,
                               gs_available=gs_available,
                               soffice_available=soffice_available)
    except Exception as e:
        logger.error(f"Error rendering index page: {e}", exc_info=True)
        return make_error_response("err-unknown", 500)

@app.route('/convert', methods=['POST'])
@limiter.limit("10 per minute")
def convert_file():
    # Giữ nguyên logic route này
    # Nó sẽ tự động sử dụng SOFFICE_PATH đã được xác định ở đầu file
    output_path = temp_libreoffice_output = input_path_for_process = None
    saved_input_paths = []; actual_conversion_type = None; start_time = time.time()
    error_key = "err-conversion"; conversion_success = False
    response_to_send = None
    try:
        if 'file' not in request.files: raise RuntimeError("err-select-file")
        file = request.files['file']
        if not file or not file.filename: raise RuntimeError("err-select-file")
        filename = secure_filename(file.filename)
        file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        allowed_office_ext = {'pdf', 'docx', 'ppt', 'pptx'}
        if not _allowed_file_extension(filename, allowed_office_ext): raise RuntimeError("err-invalid-mime-type")
        actual_conversion_type = request.form.get('conversion_type')
        valid_conversion_types = ['pdf_to_docx', 'docx_to_pdf', 'pdf_to_ppt', 'ppt_to_pdf']
        if not actual_conversion_type or actual_conversion_type not in valid_conversion_types: raise RuntimeError("err-select-conversion")
        required_ext = []
        if actual_conversion_type == 'pdf_to_docx': required_ext = ['pdf']
        elif actual_conversion_type == 'docx_to_pdf': required_ext = ['docx']
        elif actual_conversion_type == 'pdf_to_ppt': required_ext = ['pdf']
        elif actual_conversion_type == 'ppt_to_pdf': required_ext = ['ppt', 'pptx']
        if file_ext not in required_ext:
             error_key_cv = "err-format-docx" if 'docx' in required_ext else "err-format-ppt" if 'ppt' in required_ext or 'pptx' in required_ext else "err-format-pdf"
             logger.warning(f"Ext mismatch: file '{filename}' ({file_ext}), required {required_ext} for type '{actual_conversion_type}'")
             raise RuntimeError(error_key_cv)
        logger.info(f"Request /convert: file='{filename}', type='{actual_conversion_type}'")
        detected_mime = get_actual_mime_type(file) # Dùng hàm đã sửa
        if detected_mime: # Chỉ kiểm tra MIME nếu lấy được
             expected_mimes = []
             if actual_conversion_type in ['pdf_to_docx', 'pdf_to_ppt']: expected_mimes = ALLOWED_MIME_TYPES['pdf']
             elif actual_conversion_type == 'docx_to_pdf': expected_mimes = ALLOWED_MIME_TYPES['docx']
             elif actual_conversion_type == 'ppt_to_pdf': expected_mimes = ALLOWED_MIME_TYPES['ppt'] + ALLOWED_MIME_TYPES['pptx']

             if detected_mime not in expected_mimes:
                 is_expected_office_ext = file_ext in ['ppt', 'pptx', 'docx']
                 is_office_input_conversion = actual_conversion_type in ['ppt_to_pdf', 'docx_to_pdf']
                 is_pdf_input = actual_conversion_type in ['pdf_to_docx', 'pdf_to_ppt']
                 if detected_mime == 'application/octet-stream' and is_expected_office_ext and is_office_input_conversion: raise RuntimeError("err-mime-unidentified-office")
                 elif file_ext == 'pdf' and is_pdf_input and detected_mime != 'application/pdf': logger.warning(f"MIME mismatch for PDF '{filename}'. Proceeding by extension.")
                 elif detected_mime != 'application/pdf' and is_pdf_input: raise RuntimeError("err-invalid-mime-type")
                 elif not is_pdf_input: raise RuntimeError("err-invalid-mime-type")

        logger.info(f"MIME validated (or bypassed if unavailable) for {filename}: {detected_mime or 'Unavailable'}")
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        timestamp = time.strftime("%Y%m%d-%H%M%S")
        input_filename_ts = f"input_{timestamp}_{filename}"
        input_path_for_process = os.path.join(UPLOAD_FOLDER, input_filename_ts)
        try: file.seek(0); file.save(input_path_for_process); saved_input_paths.append(input_path_for_process); logger.info(f"Input saved: {input_path_for_process}")
        except Exception as save_err: logger.error(f"File save failed {filename}: {save_err}"); raise RuntimeError("err-unknown") from save_err
        base_name = filename.rsplit('.', 1)[0]
        out_ext_map = {'pdf_to_docx': 'docx', 'docx_to_pdf': 'pdf', 'pdf_to_ppt': 'pptx', 'ppt_to_pdf': 'pdf'}
        out_ext = out_ext_map.get(actual_conversion_type)
        output_filename = f"converted_{timestamp}_{secure_filename(base_name)}.{out_ext}"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)

        try:
            if actual_conversion_type == 'pdf_to_docx':
                cv = None
                try:
                    logger.info(f"Starting pdf2docx for {input_path_for_process}")
                    cv = Converter(input_path_for_process)
                    cv.convert(output_path) # No start/end arguments needed in basic usage
                    cv.close() # Quan trọng: đóng file sau khi convert
                    # Kiểm tra file output tồn tại và có size > 0
                    if os.path.isfile(output_path) and os.path.getsize(output_path) > 0:
                        conversion_success = True
                        logger.info(f"pdf2docx successful: {output_path}")
                    else:
                        logger.error(f"pdf2docx ran but output file is missing or empty: {output_path}")
                        error_key = "err-conversion" # Lỗi chung
                except (ValueError, RuntimeError, PDFPageCountError, PDFSyntaxError, Exception) as pdf2docx_err: # Bắt nhiều loại lỗi hơn
                    err_str = str(pdf2docx_err).lower()
                    if "encrypted" in err_str or "password" in err_str or "decrypt" in err_str or "err-pdf-protected" in err_str: error_key = "err-pdf-protected"
                    elif "corrupt" in err_str or "eof marker" in err_str or "invalid" in err_str or "err-pdf-corrupt" in err_str: error_key = "err-pdf-corrupt"
                    elif "no pages" in err_str or "err-pdf-no-pages" in err_str: error_key = "err-pdf-no-pages" # Thêm lỗi ko có trang
                    else: logger.error(f"pdf2docx failed: {pdf2docx_err}", exc_info=True); error_key = "err-conversion"
                finally:
                    # Đảm bảo cv được close ngay cả khi có lỗi (nếu nó đã được khởi tạo)
                    if cv:
                         try: cv.close()
                         except Exception: pass # Bỏ qua lỗi khi close nếu có
                if not conversion_success: raise RuntimeError(error_key)

            elif actual_conversion_type in ['docx_to_pdf', 'ppt_to_pdf']:
                if not SOFFICE_PATH: raise RuntimeError("err-libreoffice")
                output_dir = os.path.dirname(output_path); input_file_ext_actual = os.path.splitext(input_path_for_process)[1].lower(); expected_lo_output_name = os.path.basename(input_path_for_process).replace(input_file_ext_actual, '.pdf'); temp_libreoffice_output = os.path.join(output_dir, expected_lo_output_name); safe_remove(temp_libreoffice_output)
                cmd = [SOFFICE_PATH, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path_for_process]; logger.info(f"Running LO: {' '.join(cmd)}")
                try:
                    result = subprocess.run(cmd, check=True, timeout=LIBREOFFICE_TIMEOUT, capture_output=True, text=True, encoding='utf-8', errors='ignore')
                    logger.info(f"LO stdout:\n{result.stdout}")
                    if result.stderr: logger.warning(f"LO stderr:\n{result.stderr}")
                    if os.path.exists(temp_libreoffice_output) and os.path.getsize(temp_libreoffice_output) > 0:
                        os.rename(temp_libreoffice_output, output_path)
                        conversion_success = True
                        logger.info(f"LO conversion successful: {output_path}")
                    else:
                         # Kiểm tra xem có lỗi stderr cụ thể không
                         if result.stderr and "error" in result.stderr.lower():
                              logger.error(f"LO stderr indicates error during conversion: {result.stderr}")
                              error_key = "err-libreoffice"
                         else:
                              logger.error(f"LO ran but output '{temp_libreoffice_output}' missing/empty.")
                              error_key = "err-libreoffice"

                except subprocess.TimeoutExpired: logger.error(f"LO timed out ({LIBREOFFICE_TIMEOUT}s)."); error_key = "err-conversion-timeout"; safe_remove(temp_libreoffice_output)
                except subprocess.CalledProcessError as lo_err:
                    logger.error(f"LO failed. RC: {lo_err.returncode}")
                    if lo_err.stdout: logger.error(f"LO stdout:\n{lo_err.stdout}")
                    if lo_err.stderr: logger.error(f"LO stderr:\n{lo_err.stderr}")
                    error_key = "err-libreoffice"
                    safe_remove(temp_libreoffice_output)
                except FileNotFoundError: logger.error(f"LO not found: {SOFFICE_PATH}"); error_key = "err-libreoffice"
                except Exception as lo_run_err: logger.error(f"Unexpected LO error: {lo_run_err}", exc_info=True); error_key = "err-libreoffice"; safe_remove(temp_libreoffice_output)
                if not conversion_success: raise RuntimeError(error_key)

            elif actual_conversion_type == 'pdf_to_ppt':
                 # Sử dụng hàm đã được cải thiện
                 try:
                     if convert_pdf_to_pptx_python(input_path_for_process, output_path):
                         conversion_success = True
                         error_key = None # Thành công, xóa error key
                         logger.info("PDF->PPTX successful (Python image-based).")
                     else:
                         # Hàm trả về False, nhưng nên raise lỗi cụ thể bên trong hàm đó
                         logger.error("convert_pdf_to_pptx_python returned False without raising exception.")
                         error_key = "err-conversion" # Lỗi chung nếu hàm không raise
                 except ValueError as ve: # Bắt lỗi cụ thể từ hàm con
                      error_key = str(ve) if str(ve).startswith("err-") else "err-conversion"; logger.error(f"PDF->PPTX ValueError: {error_key}")
                 except RuntimeError as rte:
                      error_key = str(rte) if str(rte).startswith("err-") else "err-conversion"; logger.error(f"PDF->PPTX RuntimeError: {error_key}")
                 except Exception as py_ppt_err:
                      error_key = "err-unknown"; logger.error(f"Unexpected Python PDF->PPTX error: {py_ppt_err}", exc_info=True)

                 # Fallback logic giữ nguyên nhưng giờ dựa vào error_key đã được set
                 can_fallback = ( not conversion_success and SOFFICE_PATH and error_key not in ["err-pdf-corrupt", "err-pdf-protected", "err-poppler-missing"] )
                 if can_fallback:
                    logger.info(f"Python PDF->PPTX failed ({error_key}), attempting LO fallback...")
                    # Reset error key cho fallback
                    error_key_fallback = "err-libreoffice"
                    output_dir = os.path.dirname(output_path); input_file_ext_actual = os.path.splitext(input_path_for_process)[1].lower(); expected_lo_output_name = os.path.basename(input_path_for_process).replace(input_file_ext_actual, '.pptx'); temp_libreoffice_output = os.path.join(output_dir, expected_lo_output_name); safe_remove(temp_libreoffice_output)
                    cmd = [SOFFICE_PATH, '--headless', '--convert-to', 'pptx', '--outdir', output_dir, input_path_for_process]; logger.info(f"Running LO fallback: {' '.join(cmd)}")
                    try:
                        result = subprocess.run(cmd, check=True, timeout=LIBREOFFICE_TIMEOUT, capture_output=True, text=True, encoding='utf-8', errors='ignore')
                        logger.info(f"LO fallback stdout:\n{result.stdout}")
                        if result.stderr: logger.warning(f"LO fallback stderr:\n{result.stderr}")
                        if os.path.exists(temp_libreoffice_output) and os.path.getsize(temp_libreoffice_output) > 0:
                             os.rename(temp_libreoffice_output, output_path)
                             conversion_success = True
                             error_key = None # Fallback thành công
                             logger.info("LO fallback for PDF->PPTX successful.")
                        else:
                             logger.error("LO fallback ran but output missing/empty."); error_key = error_key_fallback # Giữ lỗi LO
                    except subprocess.TimeoutExpired: logger.error("LO fallback timed out."); error_key = "err-conversion-timeout"; safe_remove(temp_libreoffice_output)
                    except subprocess.CalledProcessError as lo_err:
                        logger.error(f"LO fallback failed. RC: {lo_err.returncode}")
                        if lo_err.stdout: logger.error(f"LO stdout:\n{lo_err.stdout}")
                        if lo_err.stderr: logger.error(f"LO stderr:\n{lo_err.stderr}")
                        error_key = error_key_fallback; safe_remove(temp_libreoffice_output)
                    except FileNotFoundError: logger.error(f"LO not found: {SOFFICE_PATH}"); error_key = error_key_fallback
                    except Exception as lo_run_err: logger.error(f"Unexpected LO fallback error: {lo_run_err}", exc_info=True); error_key = error_key_fallback; safe_remove(temp_libreoffice_output)
                 elif not conversion_success:
                      # Không fallback hoặc fallback không thành công, giữ lỗi ban đầu
                      logger.warning(f"Skipping or failed LO fallback. Final conversion error: {error_key}")

                 if not conversion_success:
                      # Raise lỗi cuối cùng sau khi đã thử các cách
                      raise RuntimeError(error_key or "err-conversion") # Đảm bảo luôn có error key

        except RuntimeError as rt_err: error_key = str(rt_err) if str(rt_err).startswith("err-") else "err-unknown"; logger.error(f"Caught RuntimeError during conversion: {error_key}", exc_info=False); raise
        except ValueError as val_err: error_key = str(val_err) if str(val_err).startswith("err-") else "err-unknown"; logger.error(f"Caught ValueError during conversion: {error_key}", exc_info=False); raise
        except Exception as conv_err: error_key = "err-unknown"; logger.error(f"Unexpected conversion error: {conv_err}", exc_info=True); raise

        # Gửi file nếu thành công
        if conversion_success and output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            mimetype_map = {'pdf': 'application/pdf', 'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}
            mimetype = mimetype_map.get(out_ext, 'application/octet-stream')
            try:
                response = send_file(output_path, as_attachment=True, download_name=output_filename, mimetype=mimetype)
                @response.call_on_close
                def cleanup_success():
                    logger.debug(f"Cleanup success /convert: In: {input_path_for_process}, Out: {output_path}, TempLO: {temp_libreoffice_output if 'temp_libreoffice_output' in locals() else 'N/A'}")
                    safe_remove(input_path_for_process)
                    safe_remove(output_path)
                    # Chỉ xóa temp LO nếu nó được tạo ra trong route này
                    if 'temp_libreoffice_output' in locals() and os.path.exists(temp_libreoffice_output):
                        safe_remove(temp_libreoffice_output)

                logger.info(f"Conversion successful. Sending: {output_filename}. Time: {time.time() - start_time:.2f}s")
                response_to_send = response
            except Exception as send_err:
                 logger.error(f"Error sending file {output_filename}: {send_err}", exc_info=True)
                 # Đã convert thành công nhưng ko gửi đc -> lỗi server
                 raise RuntimeError("err-unknown") from send_err
        else:
            # Nếu conversion_success=True nhưng file output có vấn đề
            if conversion_success:
                 logger.error(f"Conversion reported success but output file invalid or empty: {output_path}")
                 error_key = error_key or "err-conversion" # Sử dụng lỗi đã có hoặc lỗi chung
            # Nếu conversion_success=False thì error_key đã được set
            final_error_key = error_key or "err-conversion"
            logger.error(f"Conversion failed or output invalid. Error: {final_error_key}. Time: {time.time() - start_time:.2f}s")
            raise RuntimeError(final_error_key) # Raise lỗi cuối cùng

    except Exception as e:
         # Block xử lý lỗi chung cho toàn bộ route /convert
         final_error_key = str(e) if str(e).startswith("err-") else "err-unknown"
         status_code = 400 # Mặc định
         if final_error_key == "err-unknown": status_code = 500; logger.error(f"Unexpected /convert error: {e}", exc_info=True)
         elif final_error_key == "err-file-too-large": status_code = 413
         elif final_error_key == "err-rate-limit-exceeded": status_code = 429
         elif final_error_key == "err-csrf-invalid": status_code = 400
         elif final_error_key in ["err-pdf-protected", "err-pdf-corrupt", "err-pdf-no-pages", "err-format-docx", "err-format-ppt", "err-format-pdf", "err-invalid-mime-type", "err-mime-unidentified-office", "err-select-conversion", "err-select-file"]: status_code = 400
         elif final_error_key in ["err-libreoffice", "err-poppler-missing", "err-gs-missing"]: status_code = 503 # Service Unavailable
         elif final_error_key in ["err-conversion-timeout", "err-gs-timeout"]: status_code = 504 # Gateway Timeout
         elif final_error_key in ["err-conversion", "err-conversion-img"]: status_code = 500 # Internal server error for general conversion fails

         logger.debug(f"Cleanup failed /convert (Error: {final_error_key}).");
         [safe_remove(p) for p in saved_input_paths]
         safe_remove(output_path)
         # Chỉ xóa temp LO nếu nó được tạo ra
         if 'temp_libreoffice_output' in locals() and temp_libreoffice_output and os.path.exists(temp_libreoffice_output):
              safe_remove(temp_libreoffice_output)
         return make_error_response(final_error_key, status_code)

    # Chỉ trả về response nếu nó được tạo thành công
    if response_to_send: return response_to_send
    else:
         # Trường hợp không mong muốn: không có lỗi nhưng cũng không có response
         logger.error("Reached end of /convert without valid response or error raised.")
         return make_error_response(error_key or "err-unknown", 500)


@app.route('/convert_image', methods=['POST'])
@limiter.limit("10 per minute")
def convert_image_route():
    # Giữ nguyên logic route này
    output_path = input_path_for_pdf_input = temp_upload_dir = None
    saved_input_paths = []; actual_conversion_type = None; output_filename = None
    start_time = time.time(); error_key = "err-conversion"; conversion_success = False
    valid_files_for_processing = []; response_to_send = None
    try:
        uploaded_files = request.files.getlist('image_file')
        if not uploaded_files or not all(f and f.filename for f in uploaded_files): raise RuntimeError("err-select-file")
        logger.info(f"Request /convert_image: Received {len(uploaded_files)} file(s).")
        first_file = uploaded_files[0]; first_filename = secure_filename(first_file.filename)
        first_ext = first_filename.rsplit('.', 1)[-1].lower() if '.' in first_filename else ''
        validation_error_key = None; out_ext = None
        if first_ext == 'pdf':
            if len(uploaded_files) > 1: validation_error_key = "err-image-single-pdf"
            elif not _allowed_file_extension(first_filename, ALLOWED_IMAGE_EXTENSIONS): validation_error_key = "err-image-format"
            else:
                mime_type = get_actual_mime_type(first_file) # Dùng hàm đã sửa
                if mime_type and mime_type not in ALLOWED_MIME_TYPES['pdf']:
                    logger.warning(f"Invalid MIME for PDF {first_filename}: {mime_type}")
                    validation_error_key = "err-invalid-mime-type"
                elif not mime_type and magic: # Nếu magic có nhưng ko detect đc
                     logger.warning(f"Could not detect MIME for PDF {first_filename}. Proceeding by extension.")
                     actual_conversion_type = 'pdf_to_image'; out_ext = 'zip'; valid_files_for_processing.append(first_file)
                else: # Hoặc magic không có, hoặc detect đúng
                     actual_conversion_type = 'pdf_to_image'; out_ext = 'zip'; valid_files_for_processing.append(first_file)

        elif first_ext in ['jpg', 'jpeg']:
            actual_conversion_type = 'image_to_pdf'; out_ext = 'pdf'; allowed_image_mimes = ALLOWED_MIME_TYPES['jpeg']
            try: temp_upload_dir = tempfile.mkdtemp(prefix="img2pdf_")
            except Exception as temp_err: logger.error(f"Failed create temp dir: {temp_err}"); raise RuntimeError("err-unknown") from temp_err
            total_size = 0; max_size_bytes = app.config['MAX_CONTENT_LENGTH']
            for i, f in enumerate(uploaded_files):
                fname_sec = secure_filename(f.filename); f_ext = fname_sec.rsplit('.', 1)[-1].lower() if '.' in fname_sec else ''
                if f_ext not in ['jpg', 'jpeg']: validation_error_key = "err-image-all-images"; logger.warning(f"Invalid ext {fname_sec}"); break
                # Kiểm tra size trước khi lưu
                f.stream.seek(0, os.SEEK_END); file_size = f.stream.tell(); f.stream.seek(0); total_size += file_size
                if total_size > max_size_bytes: validation_error_key = "err-file-too-large"; logger.warning(f"Total size limit exceeded at {fname_sec}"); break
                # Kiểm tra MIME trước khi lưu (nếu magic có)
                if magic:
                     mime_type = get_actual_mime_type(f)
                     if not mime_type or mime_type not in allowed_image_mimes:
                         validation_error_key = "err-invalid-mime-type-image"; logger.warning(f"Invalid MIME for img {fname_sec}: {mime_type}"); break
                # Lưu file vào thư mục tạm
                temp_image_path = os.path.join(temp_upload_dir, f"{i}_{fname_sec}") # Đặt tên file tạm khác nhau
                try:
                    f.save(temp_image_path)
                    valid_files_for_processing.append(temp_image_path)
                    saved_input_paths.append(temp_image_path) # Chỉ thêm vào đây nếu save thành công
                except Exception as save_err:
                    logger.error(f"Failed save temp image {fname_sec}: {save_err}")
                    validation_error_key = "err-unknown"; break
            # Kiểm tra lại sau vòng lặp
            if validation_error_key: pass # Đã có lỗi, sẽ raise sau
            elif not valid_files_for_processing: validation_error_key = "err-select-file" # Không có file nào hợp lệ được xử lý

        else: validation_error_key = "err-image-format"

        # Raise lỗi validation nếu có
        if validation_error_key: raise RuntimeError(validation_error_key)

        logger.info(f"Conversion type: {actual_conversion_type}. Validated inputs.")
        os.makedirs(UPLOAD_FOLDER, exist_ok=True); timestamp = time.strftime("%Y%m%d-%H%M%S")

        # Lưu input PDF nếu cần (chỉ cho pdf_to_image)
        if actual_conversion_type == 'pdf_to_image':
            pdf_file_storage = valid_files_for_processing[0] # Đây là FileStorage object
            input_filename_ts = f"input_{timestamp}_{secure_filename(pdf_file_storage.filename)}"
            input_path_for_pdf_input = os.path.join(UPLOAD_FOLDER, input_filename_ts)
            try:
                pdf_file_storage.stream.seek(0) # Đảm bảo đọc từ đầu
                pdf_file_storage.save(input_path_for_pdf_input)
                saved_input_paths.append(input_path_for_pdf_input) # Thêm path file đã lưu vào cleanup list
                logger.info(f"Input PDF saved: {input_path_for_pdf_input}")
            except Exception as save_err: logger.error(f"Failed save PDF input: {save_err}"); raise RuntimeError("err-unknown") from save_err

        # Xác định tên file output
        base_name = first_filename.rsplit('.', 1)[0]
        output_filename = f"converted_{timestamp}_{secure_filename(base_name)}.{out_ext}"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)

        # Thực hiện convert
        try:
            if actual_conversion_type == 'pdf_to_image':
                 conversion_success = convert_pdf_to_image_zip(input_path_for_pdf_input, output_path)
            elif actual_conversion_type == 'image_to_pdf':
                 # valid_files_for_processing bây giờ chứa các đường dẫn file ảnh đã lưu
                 conversion_success = convert_images_to_pdf(valid_files_for_processing, output_path)

            if not conversion_success:
                 # Hàm convert nên raise lỗi cụ thể, nhưng phòng trường hợp nó trả về False
                 logger.error(f"Conversion function for {actual_conversion_type} returned False.")
                 raise RuntimeError(error_key or "err-conversion") # Sử dụng lỗi đã có hoặc lỗi chung

        except ValueError as val_err: error_key = str(val_err) if str(val_err).startswith("err-") else "err-conversion"; logger.error(f"Image conversion ValueError: {error_key}", exc_info=False); raise
        except RuntimeError as rt_err: error_key = str(rt_err) if str(rt_err).startswith("err-") else "err-conversion"; logger.error(f"Image conversion RuntimeError: {error_key}", exc_info=False); raise
        except Exception as conv_err: error_key = "err-unknown"; logger.error(f"Unexpected image conversion error: {conv_err}", exc_info=True); raise

        # Gửi file nếu thành công
        if conversion_success and output_path and os.path.exists(output_path):
             # Kiểm tra size file output > 0 (tránh gửi file rỗng)
             if os.path.getsize(output_path) > 0:
                 mimetype = 'application/zip' if out_ext == 'zip' else 'application/pdf'
                 try:
                     response = send_file(output_path, as_attachment=True, download_name=output_filename, mimetype=mimetype)
                     @response.call_on_close
                     def cleanup_image_success():
                         logger.debug(f"Cleanup success /convert_image: Inputs: {saved_input_paths}, Out: {output_path}, TempDir: {temp_upload_dir}")
                         # Xóa các file input đã lưu (PDF hoặc ảnh tạm)
                         [safe_remove(p) for p in saved_input_paths]
                         safe_remove(output_path) # Xóa file output
                         safe_remove(temp_upload_dir) # Xóa thư mục tạm chứa ảnh (nếu có)

                     logger.info(f"Image conversion successful. Sending: {output_filename}. Time: {time.time() - start_time:.2f}s")
                     response_to_send = response
                 except Exception as send_err:
                      logger.error(f"Error sending image file {output_filename}: {send_err}", exc_info=True)
                      raise RuntimeError("err-unknown") from send_err
             else:
                  # File output rỗng
                  logger.error(f"Image conversion resulted in an empty output file: {output_path}")
                  final_error_key = error_key or "err-conversion"
                  raise RuntimeError(final_error_key)
        else:
             # conversion_success = False hoặc file output không tồn tại
             final_error_key = error_key or "err-conversion"
             logger.error(f"Image conversion failed or output invalid. Error: {final_error_key}. Time: {time.time() - start_time:.2f}s")
             raise RuntimeError(final_error_key)

    except Exception as e:
        # Block xử lý lỗi chung cho /convert_image
        final_error_key = str(e) if str(e).startswith("err-") else "err-unknown"; status_code = 400
        if final_error_key == "err-unknown": status_code = 500; logger.error(f"Unexpected /convert_image error: {e}", exc_info=True)
        elif final_error_key == "err-file-too-large": status_code = 413
        elif final_error_key == "err-rate-limit-exceeded": status_code = 429
        elif final_error_key == "err-csrf-invalid": status_code = 400
        elif final_error_key in ["err-pdf-protected", "err-pdf-corrupt", "err-invalid-image-file", "err-image-format", "err-image-single-pdf", "err-image-all-images", "err-invalid-mime-type", "err-invalid-mime-type-image", "err-select-file"]: status_code = 400
        elif final_error_key in ["err-poppler-missing"]: status_code = 503
        elif final_error_key in ["err-conversion", "err-conversion-img", "err-poppler-check-failed"]: status_code = 500

        logger.debug(f"Cleanup failed /convert_image (Error: {final_error_key}).")
        # Dọn dẹp kỹ hơn: xóa cả input đã lưu và thư mục tạm
        [safe_remove(p) for p in saved_input_paths]
        safe_remove(output_path)
        safe_remove(temp_upload_dir)
        return make_error_response(final_error_key, status_code)

    # Chỉ trả về nếu có response hợp lệ
    if response_to_send: return response_to_send
    else:
         logger.error("Reached end of /convert_image without valid response or error raised.")
         return make_error_response(error_key or "err-unknown", 500)


@app.route('/compress_pdf', methods=['POST'])
@limiter.limit("10 per minute")
def compress_pdf_route():
    # Giữ nguyên logic route này
    # Nó sẽ tự động sử dụng GS_PATH đã được xác định ở đầu file
    input_path = output_path = None; saved_input_paths = []
    start_time = time.time(); error_key = "err-gs-failed"; compression_success = False
    response_to_send = None
    try:
        if not GS_PATH: raise RuntimeError("err-gs-missing")
        if 'file' not in request.files: raise RuntimeError("err-select-file")
        file = request.files['file']
        if not file or not file.filename: raise RuntimeError("err-select-file")
        filename = secure_filename(file.filename); file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        if file_ext != 'pdf': logger.warning(f"Rejected non-PDF for compression: {filename}"); raise RuntimeError("err-format-pdf")

        detected_mime = get_actual_mime_type(file) # Dùng hàm đã sửa
        if detected_mime and detected_mime not in ALLOWED_MIME_TYPES['pdf']:
            logger.warning(f"MIME check failed for compression {filename}: '{detected_mime}'")
            raise RuntimeError("err-invalid-mime-type")
        elif not detected_mime and magic: # Nếu magic có nhưng ko detect đc
            logger.warning(f"Could not detect MIME for compression {filename}. Proceeding by extension.")
        # Bỏ qua kiểm tra MIME nếu magic không có

        quality = request.form.get('quality', 'medium')
        if quality not in ['low', 'medium', 'high']: logger.warning(f"Invalid quality level specified: {quality}"); raise RuntimeError("err-invalid-quality")
        logger.info(f"Request /compress_pdf: file='{filename}', quality='{quality}'")

        os.makedirs(UPLOAD_FOLDER, exist_ok=True); timestamp = time.strftime("%Y%m%d-%H%M%S"); input_filename_ts = f"input_{timestamp}_{filename}"; input_path = os.path.join(UPLOAD_FOLDER, input_filename_ts)
        try:
             file.seek(0); file.save(input_path); saved_input_paths.append(input_path); logger.info(f"Input PDF saved for compression: {input_path}")
        except Exception as save_err: logger.error(f"Save failed for compression {filename}: {save_err}"); raise RuntimeError("err-unknown") from save_err

        base_name = filename.rsplit('.', 1)[0]; output_filename_base = secure_filename(f"{base_name}_compressed_{quality}"); output_filename = f"{output_filename_base}.pdf"; output_path = os.path.join(UPLOAD_FOLDER, output_filename)

        try:
            # Gọi hàm compress đã được cải thiện
            compression_success = compress_pdf_ghostscript(input_path, output_path, quality)
            if not compression_success:
                 # Hàm compress nên raise lỗi, nhưng phòng trường hợp nó trả về False
                 raise RuntimeError(error_key or "err-gs-failed") # Sử dụng lỗi đã có hoặc mặc định
        except ValueError as ve: error_key = str(ve) if str(ve).startswith("err-") else "err-gs-failed"; logger.warning(f"Compression ValueError: {error_key}"); raise
        except RuntimeError as rt_err: error_key = str(rt_err) if str(rt_err).startswith("err-") else "err-gs-failed"; logger.error(f"Caught RuntimeError compression: {error_key}", exc_info=False); raise
        except Exception as comp_err: error_key = "err-unknown"; logger.error(f"Unexpected compression error: {comp_err}", exc_info=True); raise

        # Gửi file nếu thành công và file output hợp lệ
        if compression_success and output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            try:
                response = send_file(output_path, as_attachment=True, download_name=output_filename, mimetype='application/pdf')
                @response.call_on_close
                def cleanup_compress_success():
                    logger.debug(f"Cleanup success /compress_pdf: In: {input_path}, Out: {output_path}")
                    safe_remove(input_path)
                    safe_remove(output_path)
                logger.info(f"Compression successful. Sending: {output_filename}. Time: {time.time() - start_time:.2f}s")
                response_to_send = response
            except Exception as send_err:
                 logger.error(f"Error sending compressed file {output_filename}: {send_err}", exc_info=True)
                 raise RuntimeError("err-unknown") from send_err
        else:
             if compression_success: # Thành công nhưng output không hợp lệ
                  logger.error(f"Compression reported success but output file invalid or empty: {output_path}")
                  error_key = error_key or "err-gs-failed"
             # Nếu compression_success=False thì error_key đã được set
             final_error_key = error_key
             logger.error(f"Compression failed or output invalid. Error: {final_error_key}. Time: {time.time() - start_time:.2f}s")
             raise RuntimeError(final_error_key)

    except Exception as e:
        # Block xử lý lỗi chung cho /compress_pdf
        final_error_key = str(e) if str(e).startswith("err-") else "err-unknown"; status_code = 400
        if final_error_key == "err-unknown": status_code = 500; logger.error(f"Unexpected /compress_pdf error: {e}", exc_info=True)
        elif final_error_key == "err-file-too-large": status_code = 413
        elif final_error_key == "err-rate-limit-exceeded": status_code = 429
        elif final_error_key == "err-csrf-invalid": status_code = 400
        elif final_error_key in ["err-pdf-protected", "err-pdf-corrupt", "err-format-pdf", "err-invalid-mime-type", "err-invalid-quality", "err-select-file"]: status_code = 400
        elif final_error_key in ["err-gs-failed", "err-gs-missing"]: status_code = 503
        elif final_error_key == "err-gs-timeout": status_code = 504
        # Phân loại lỗi conversion chung
        elif final_error_key == "err-conversion": status_code = 500

        logger.debug(f"Cleanup failed /compress_pdf (Error: {final_error_key}).")
        [safe_remove(p) for p in saved_input_paths]
        safe_remove(output_path)
        return make_error_response(final_error_key, status_code)

    if response_to_send: return response_to_send
    else:
         logger.error("Reached end of /compress_pdf without valid response or error raised.")
         return make_error_response(error_key or "err-unknown", 500)

@app.route('/compress_docx', methods=['POST'])
@limiter.limit("10 per minute")
def compress_docx_route():
    # Giữ nguyên logic route này
    # Nó sẽ tự động sử dụng SOFFICE_PATH và GS_PATH đã được xác định
    input_path_docx = temp_pdf_uncompressed = temp_pdf_compressed = final_output_docx = None
    saved_input_paths = []; intermediate_files = []
    start_time = time.time(); error_key = "err-conversion"; process_success = False
    response_to_send = None
    final_output_filename_base = None
    final_download_name = None
    try:
        if not SOFFICE_PATH: raise RuntimeError("err-libreoffice")
        if not GS_PATH: raise RuntimeError("err-gs-missing")
        if 'file' not in request.files: raise RuntimeError("err-select-file")
        file = request.files['file']
        if not file or not file.filename: raise RuntimeError("err-select-file")
        filename = secure_filename(file.filename); file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        if file_ext != 'docx': logger.warning(f"Rejected non-DOCX: {filename}"); raise RuntimeError("err-format-docx")

        detected_mime = get_actual_mime_type(file) # Dùng hàm đã sửa
        if detected_mime: # Chỉ kiểm tra nếu có
            if detected_mime not in ALLOWED_MIME_TYPES['docx']:
                if detected_mime == 'application/octet-stream':
                    logger.warning(f"Unidentified MIME for DOCX {filename}. Proceeding with caution.");
                    # Có thể raise lỗi ở đây nếu muốn chặt chẽ hơn
                    # raise RuntimeError("err-mime-unidentified-office")
                else:
                     logger.warning(f"MIME check failed for DOCX {filename}: '{detected_mime}'")
                     raise RuntimeError("err-invalid-mime-type")
        # Bỏ qua nếu magic không có

        logger.info(f"Request /compress_docx: file='{filename}'")
        os.makedirs(UPLOAD_FOLDER, exist_ok=True); timestamp = time.strftime("%Y%m%d-%H%M%S"); input_filename_ts = f"input_{timestamp}_{filename}"; input_path_docx = os.path.join(UPLOAD_FOLDER, input_filename_ts)
        try: file.seek(0); file.save(input_path_docx); saved_input_paths.append(input_path_docx); logger.info(f"Input DOCX saved: {input_path_docx}")
        except Exception as save_err: logger.error(f"Save failed for DOCX {filename}: {save_err}"); raise RuntimeError("err-unknown") from save_err
        base_name = filename.rsplit('.', 1)[0]; output_dir = os.path.dirname(input_path_docx)

        # Đặt tên file tạm và file cuối cùng
        pdf_base_name = os.path.basename(input_path_docx).replace('.docx', '.pdf')
        temp_pdf_uncompressed = os.path.join(output_dir, f"temp_uncomp_{timestamp}_{pdf_base_name}"); intermediate_files.append(temp_pdf_uncompressed)
        temp_pdf_compressed = os.path.join(output_dir, f"temp_comp_{timestamp}_{pdf_base_name}"); intermediate_files.append(temp_pdf_compressed)
        final_output_filename_base = secure_filename(f"{base_name}_compressed")
        final_output_docx = os.path.join(UPLOAD_FOLDER, f"{final_output_filename_base}.docx")
        final_download_name = f"{final_output_filename_base}.docx"

        # --- Step 1: Convert DOCX to Uncompressed PDF ---
        lo_success = False; lo_direct_output_path = None # Khởi tạo để cleanup
        try:
            cmd_lo = [SOFFICE_PATH, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path_docx]
            # Tạo tên output mong đợi của LO để rename sau
            lo_direct_output_path = os.path.join(output_dir, pdf_base_name)
            safe_remove(lo_direct_output_path) # Xóa nếu có từ lần chạy trước
            logger.info(f"Running LO for DOCX->PDF: {' '.join(cmd_lo)}")
            result_lo = subprocess.run(cmd_lo, check=True, timeout=LIBREOFFICE_TIMEOUT, capture_output=True, text=True, encoding='utf-8', errors='ignore')
            logger.info(f"LO stdout (DOCX->PDF):\n{result_lo.stdout}")
            if result_lo.stderr: logger.warning(f"LO stderr (DOCX->PDF):\n{result_lo.stderr}")

            if os.path.exists(lo_direct_output_path) and os.path.getsize(lo_direct_output_path) > 0:
                os.rename(lo_direct_output_path, temp_pdf_uncompressed) # Đổi tên thành file tạm của chúng ta
                lo_success = True
                logger.info(f"LO DOCX->PDF successful: {temp_pdf_uncompressed}")
            else:
                 logger.error(f"LO ran but expected PDF output '{lo_direct_output_path}' missing/empty.")
                 error_key = "err-libreoffice"
                 # Không cần xóa lo_direct_output_path vì nó không tồn tại hoặc rỗng
        except subprocess.TimeoutExpired: logger.error(f"LO DOCX->PDF timed out ({LIBREOFFICE_TIMEOUT}s)."); error_key = "err-conversion-timeout"; safe_remove(lo_direct_output_path)
        except subprocess.CalledProcessError as lo_err:
            logger.error(f"LO DOCX->PDF failed. RC: {lo_err.returncode}")
            if lo_err.stdout: logger.error(f"LO stdout:\n{lo_err.stdout}")
            if lo_err.stderr: logger.error(f"LO stderr:\n{lo_err.stderr}")
            error_key = "err-libreoffice"; safe_remove(lo_direct_output_path)
        except FileNotFoundError: logger.error(f"LO not found: {SOFFICE_PATH}"); error_key = "err-libreoffice"
        except Exception as lo_run_err: logger.error(f"Unexpected LO error DOCX->PDF: {lo_run_err}", exc_info=True); error_key = "err-libreoffice"; safe_remove(lo_direct_output_path)
        if not lo_success: raise RuntimeError(error_key)

        # --- Step 2: Compress the intermediate PDF ---
        gs_success = False
        if os.path.exists(temp_pdf_uncompressed) and os.path.getsize(temp_pdf_uncompressed) > 0:
            try:
                # Sử dụng hàm compress đã sửa, chất lượng 'low' để nén tối đa cho DOCX
                gs_success = compress_pdf_ghostscript(temp_pdf_uncompressed, temp_pdf_compressed, quality_level='low')
                if not gs_success: raise RuntimeError(error_key or "err-gs-failed") # Sử dụng lỗi đã có nếu compress_pdf_ghostscript raise
            except ValueError as ve: error_key = str(ve) if str(ve).startswith("err-") else "err-gs-failed"; logger.warning(f"GS compression ValueError: {error_key}"); raise
            except RuntimeError as rt_err: error_key = str(rt_err) if str(rt_err).startswith("err-") else "err-gs-failed"; logger.error(f"Caught RuntimeError GS compress: {error_key}", exc_info=False); raise
            except Exception as comp_err: error_key = "err-unknown"; logger.error(f"Unexpected GS compress error: {comp_err}", exc_info=True); raise
        else:
             logger.error(f"Intermediate uncompressed PDF '{temp_pdf_uncompressed}' missing/empty after LO step.")
             raise RuntimeError("err-libreoffice") # Lỗi từ bước trước
        # if not gs_success: raise RuntimeError(error_key) # Không cần nữa vì các exception đã được raise

        # --- Step 3: Convert Compressed PDF back to DOCX ---
        pdf2docx_success = False
        if os.path.exists(temp_pdf_compressed) and os.path.getsize(temp_pdf_compressed) > 0:
            cv = None
            try:
                logger.info(f"Starting pdf2docx for compressed PDF {temp_pdf_compressed}")
                cv = Converter(temp_pdf_compressed)
                cv.convert(final_output_docx)
                cv.close() # Đóng file
                # Kiểm tra file cuối cùng
                if os.path.exists(final_output_docx) and os.path.getsize(final_output_docx) > 0:
                    pdf2docx_success = True
                    logger.info(f"pdf2docx conversion successful: {final_output_docx}")
                else:
                     logger.error(f"pdf2docx ran but final DOCX file is missing or empty: {final_output_docx}")
                     error_key = "err-conversion" # Lỗi chung
            except (ValueError, RuntimeError, PDFPageCountError, PDFSyntaxError, Exception) as pdf2docx_err: # Bắt lỗi rộng hơn
                err_str = str(pdf2docx_err).lower()
                if "encrypted" in err_str or "password" in err_str or "decrypt" in err_str or "err-pdf-protected" in err_str: error_key = "err-pdf-protected"
                elif "corrupt" in err_str or "eof marker" in err_str or "invalid" in err_str or "err-pdf-corrupt" in err_str: error_key = "err-pdf-corrupt"
                elif "no pages" in err_str or "err-pdf-no-pages" in err_str: error_key = "err-pdf-no-pages"
                else: logger.error(f"pdf2docx conversion failed for {temp_pdf_compressed}: {pdf2docx_err}", exc_info=True); error_key = "err-conversion"
            finally:
                if cv:
                     try: cv.close()
                     except Exception: pass
        else:
             logger.error(f"Intermediate compressed PDF '{temp_pdf_compressed}' missing/empty after GS step.")
             error_key = "err-gs-failed" # Lỗi từ bước trước
        if not pdf2docx_success: raise RuntimeError(error_key)

        # --- Step 4: Handle Success ---
        process_success = True # Đã vượt qua các bước
        if final_output_docx and os.path.exists(final_output_docx) and os.path.getsize(final_output_docx) > 0:
            try:
                final_mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                response = send_file(final_output_docx, as_attachment=True, download_name=final_download_name, mimetype=final_mimetype)
                @response.call_on_close
                def cleanup_compress_docx_success():
                    logger.debug(f"Cleanup success /compress_docx: Input: {input_path_docx}, TempUncomp: {temp_pdf_uncompressed}, TempComp: {temp_pdf_compressed}, Final: {final_output_docx}")
                    safe_remove(input_path_docx)
                    # Xóa các file trung gian
                    [safe_remove(f) for f in intermediate_files]
                    safe_remove(final_output_docx)
                logger.info(f"DOCX compression successful. Sending: {final_download_name}. Time: {time.time() - start_time:.2f}s")
                response_to_send = response
            except Exception as send_err:
                 logger.error(f"Error sending final DOCX file {final_download_name}: {send_err}", exc_info=True)
                 raise RuntimeError("err-unknown") from send_err
        else:
             # Process thành công nhưng file cuối cùng lỗi?
             logger.error(f"Process finished but final DOCX file path '{final_output_docx}' missing/empty.")
             raise RuntimeError("err-conversion") # Lỗi chung

    except Exception as e:
        # Block xử lý lỗi chung cho /compress_docx
        final_error_key = str(e) if str(e).startswith("err-") else "err-unknown"; status_code = 400
        if final_error_key == "err-unknown": status_code = 500; logger.error(f"Unexpected /compress_docx error: {e}", exc_info=True)
        elif final_error_key == "err-file-too-large": status_code = 413
        elif final_error_key == "err-rate-limit-exceeded": status_code = 429
        elif final_error_key == "err-csrf-invalid": status_code = 400
        elif final_error_key in ["err-format-docx", "err-invalid-mime-type", "err-mime-unidentified-office", "err-select-file"]: status_code = 400
        elif final_error_key in ["err-pdf-protected", "err-pdf-corrupt", "err-pdf-no-pages"]: status_code = 400 # Lỗi có thể xảy ra ở bước pdf2docx cuối
        elif final_error_key in ["err-libreoffice", "err-gs-missing", "err-gs-failed", "err-conversion"]: status_code = 503 # Coi lỗi engine là Service Unavail
        elif final_error_key in ["err-conversion-timeout", "err-gs-timeout"]: status_code = 504

        logger.debug(f"Cleanup failed /compress_docx (Error: {final_error_key}).")
        [safe_remove(p) for p in saved_input_paths]
        [safe_remove(f) for f in intermediate_files] # Đảm bảo xóa file trung gian khi lỗi
        safe_remove(final_output_docx)
        # Xóa cả output trực tiếp của LO nếu nó được tạo ra và còn tồn tại
        if 'lo_direct_output_path' in locals() and lo_direct_output_path and os.path.exists(lo_direct_output_path):
             safe_remove(lo_direct_output_path)
        return make_error_response(final_error_key, status_code)

    if response_to_send: return response_to_send
    else:
        logger.error("Reached end of /compress_docx without valid response or error raised.")
        return make_error_response(error_key or "err-unknown", 500)

@app.teardown_appcontext
def cleanup_old_files(exception=None):
    # Giữ nguyên logic dọn dẹp này
    if not os.path.exists(UPLOAD_FOLDER): return
    logger.debug("Running teardown cleanup...")
    try:
        now = time.time(); max_age = 3600 # 1 hour
        deleted_count = 0; checked_count = 0
        try:
            items = os.listdir(UPLOAD_FOLDER)
        except OSError as list_err:
             logger.error(f"Teardown listdir error for {UPLOAD_FOLDER}: {list_err}")
             return
        for item_name in items:
            # Bỏ qua các thư mục tạm đặc biệt (có thể bạn không dùng nữa)
            # if item_name and (item_name.startswith("img2pdf_") or item_name.startswith("pdfimg_") or item_name.startswith("pdf2imgzip_")):
            #      logger.debug(f"Teardown skipping special temp dir/file: {item_name}")
            #      continue
            # Chỉ xóa file (không xóa thư mục) trong uploads
            path = os.path.join(UPLOAD_FOLDER, item_name)
            try:
                 if os.path.isfile(path): # Chỉ xóa file
                     stat_result = os.stat(path)
                     file_age = now - stat_result.st_mtime
                     checked_count += 1
                     if file_age > max_age:
                         logger.info(f"Teardown: Removing old file ({file_age:.0f}s > {max_age}s): {path}")
                         if safe_remove(path):
                              deleted_count += 1
            except FileNotFoundError: continue # File đã bị xóa bởi process khác
            except Exception as e: logger.warning(f"Teardown check error for {path}: {e}")

        if checked_count > 0 or deleted_count > 0:
             logger.info(f"Teardown cleanup: Checked {checked_count} files, removed {deleted_count} files older than {max_age} seconds in {UPLOAD_FOLDER}.")
        else: logger.debug("Teardown cleanup: No old files found/removed.")
    except Exception as e: logger.error(f"Teardown critical error: {e}", exc_info=True)

if __name__ == '__main__':
    try:
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        logger.info(f"Upload folder created/exists: {os.path.abspath(UPLOAD_FOLDER)}")
    except OSError as mkdir_err:
        logger.critical(f"FATAL: Cannot create upload folder {UPLOAD_FOLDER}: {mkdir_err}.")
        sys.exit(1)

    # Log các đường dẫn đã xác định
    logger.info(f"LibreOffice Path Used: {SOFFICE_PATH if SOFFICE_PATH else 'Not Found/Verified'}")
    logger.info(f"Ghostscript Path Used: {GS_PATH if GS_PATH else 'Not Found/Verified'}")

    csrf_enabled = app.config.get('WTF_CSRF_ENABLED', True)
    logger.info(f"CSRF Protection Enabled: {csrf_enabled}")
    logger.info(f"Rate Limiting Enabled: Yes (Default limits active)")

    port = int(os.environ.get('PORT', 5003))
    host = os.environ.get('HOST', '0.0.0.0') # Nghe trên tất cả các interface
    # Kiểm tra biến môi trường FLASK_DEBUG hoặc một biến khác nếu muốn
    # Sử dụng giá trị từ Run Configuration nếu chạy qua PyCharm
    run_debug_mode = os.environ.get('FLASK_DEBUG', 'False').lower() in ['true', '1', 't']

    logger.info(f"Starting server on {host}:{port} - Debug: {run_debug_mode}")

    if run_debug_mode:
        logger.warning("Running in Flask DEVELOPMENT mode (Debug=True, Werkzeug server).")
        # Khi debug=True, app.run() sử dụng server Werkzeug với reloader
        # Không cần waitress ở đây
        app.run(host=host, port=port, debug=True, threaded=True) # use_reloader mặc định là True khi debug=True
    else:
        logger.info("Running in PRODUCTION mode (Debug=False, Waitress server).")
        try:
            from waitress import serve
            serve(app, host=host, port=port, threads=4) # Chạy Waitress cho production
        except ImportError:
            logger.critical("Waitress not found! Install waitress for production.")
            logger.warning("FALLING BACK TO FLASK DEVELOPMENT SERVER (Werkzeug) WITHOUT DEBUG.")
            # Fallback về server dev của Flask nhưng tắt debug nếu waitress không có
            app.run(host=host, port=port, debug=False, threaded=True)